from pathlib import Path
import base64
import math
import re

import pandas as pd
import plotly.graph_objects as go
import streamlit as st
from io import BytesIO
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except Exception:
    Presentation = None
    Inches = None
    Pt = None
    RGBColor = None





def image_file_to_data_uri(path):
    """Return a data URI so local images can be rendered inside raw HTML."""
    image_path = Path(path)
    if not image_path.exists():
        return ""
    suffix = image_path.suffix.lower()
    mime = "image/png" if suffix == ".png" else "image/jpeg" if suffix in [".jpg", ".jpeg"] else "image/svg+xml"
    encoded = base64.b64encode(image_path.read_bytes()).decode("utf-8")
    return f"data:{mime};base64,{encoded}"


# =========================
# App config
# =========================

st.set_page_config(
    page_title="OEM Macro Conversion Intelligence Report",
    page_icon="🚗",
    layout="wide",
)

DATA_FILE = Path(__file__).parent / "OEM Visit to Sales Data 2024-2026.xlsx"

VALTECH_BLUE = "#003CB3"
VALTECH_GREY = "#6F6F6F"
VALTECH_LIGHT_GREY = "#F3F3F3"
BLACK = "#000000"
WHITE = "#FFFFFF"
GREEN = "#12C76B"
PINK = "#FF4B55"
AMBER = "#FFB000"
INTELLIGENCE = "#2563EB"

TOYOTA_LOGO = ""  # legacy compatibility only
LEXUS_LOGO = ""  # legacy compatibility only
VALTECH_LOGO = "https://mma.prnewswire.com/media/2728124/Valtech_Logo.jpg"

MARKETS = ["MM5", "UK", "France", "Germany", "Italy", "Spain"]

COMPARISON_OPTIONS = {
    "2024 vs 2025": {
        "previous_period": "2024",
        "current_period": "2025",
        "previous_label": "2024",
        "current_label": "2025",
        "coverage_label": "2024 & 2025",
    },
    "Jan-Apr 2025 vs Jan-Apr 2026": {
        "previous_period": "Jan - April 25",
        "current_period": "Jan - April 26",
        "previous_label": "Jan-Apr 2025",
        "current_label": "Jan-Apr 2026",
        "coverage_label": "2026 (Jan to April)",
    },
}

ACTIVE_COMPARISON = COMPARISON_OPTIONS["2024 vs 2025"]
PREVIOUS_LABEL = ACTIVE_COMPARISON["previous_label"]
CURRENT_LABEL = ACTIVE_COMPARISON["current_label"]


TOYOTA_SET = []  # legacy compatibility only; no OEM is locked or privileged
LEXUS_SET = []  # legacy compatibility only; no OEM is locked or privileged
CHINESE_SET = ["BYD Auto", "MG", "XPeng", "NIO", "Geely", "Omoda", "Jaecoo", "Leapmotor", "GWM Ora"]

OEM_CLUSTERS = {
    "Volume Leaders": [
        "BMW", "Toyota", "Mercedes-Benz", "Audi", "Ford", "Kia", "VW", "Volkswagen",
        "Nissan", "Skoda", "Honda"
    ],
    "EV Challengers": [
        "Tesla", "BYD", "BYD Auto", "Polestar", "NIO", "Zeekr", "Leapmotor", "AION"
    ],
    "European Mass Market": [
        "Vauxhall", "Opel", "Renault", "Peugeot", "Dacia", "Citroën", "Citroen",
        "SEAT", "Fiat", "DS Automobiles", "Alpine"
    ],
    "Luxury & Performance": [
        "Volvo", "Land Rover", "Porsche", "Lexus", "Ferrari", "Jaguar", "Aston Martin",
        "Bentley", "Lamborghini", "Maserati", "Rolls-Royce"
    ],
    "Asian Mainstream": [
        "Hyundai", "MINI", "Mazda", "CUPRA", "Suzuki", "Genesis", "Subaru", "Mitsubishi"
    ],
    "Chinese New Entrants": [
        "MG", "Smart", "Jaecoo", "KGM", "Omoda", "GWM", "GWM Ora", "Isuzu",
        "Chery", "SAIC Maxus", "XPeng", "Geely"
    ],
}

CLUSTER_DESCRIPTIONS = {
    "Volume Leaders": "Established mass-market brands with the largest digital footprint",
    "EV Challengers": "Electric-first brands disrupting the traditional OEM landscape",
    "European Mass Market": "European mainstream brands competing on value and heritage",
    "Luxury & Performance": "Premium and ultra-premium brands targeting high-intent affluent buyers",
    "Asian Mainstream": "Japanese and Korean mainstream brands with strong hybrid / EV portfolios",
    "Chinese New Entrants": "Chinese OEMs building digital presence ahead of market entry or expansion",
}

CLUSTER_COLORS = {
    "Volume Leaders": "#2563EB",
    "EV Challengers": "#12C76B",
    "European Mass Market": "#FFB000",
    "Luxury & Performance": "#B544F4",
    "Asian Mainstream": "#FF7A1A",
    "Chinese New Entrants": "#FF3B77",
}


LOGO_DOMAIN_MAP = {
    "Toyota": "toyota.com",
    "Lexus": "lexus.com",
    "VW": "volkswagen.com",
    "BMW": "bmw.com",
    "Mercedes-Benz": "mercedes-benz.com",
    "Audi": "audi.com",
    "Volvo": "volvocars.com",
    "Tesla": "tesla.com",
    "Ford": "ford.com",
    "Peugeot": "peugeot.com",
    "Renault": "renault.com",
    "Hyundai": "hyundai.com",
    "Kia": "kia.com",
    "Nissan": "nissan-global.com",
    "Skoda": "skoda-auto.com",
    "SEAT": "seat.com",
    "Dacia": "dacia.com",
    "Jaguar": "jaguar.com",
    "Land Rover": "landrover.com",
    "BYD Auto": "byd.com",
    "MG": "mg.co.uk",
    "XPeng": "xpeng.com",
    "Geely": "geely.com",
    "Omoda": "omodaauto.com",
    "Jaecoo": "jaecoo.com",
    "Leapmotor": "leapmotor.com",
    "GWM Ora": "gwm-global.com",
    "Porsche": "porsche.com",
    "Polestar": "polestar.com",
}


# =========================
# Styling
# =========================

st.markdown(
    """
<style>
html, body, [class*="css"] {
    font-family: "Helvetica Neue", Arial, sans-serif;
}
.stApp {
    background: #ffffff;
}
section[data-testid="stSidebar"] {
    background: #000000;
}
section[data-testid="stSidebar"] * {
    color: #ffffff;
}
section[data-testid="stSidebar"] .stMultiSelect [data-baseweb="select"] *,
section[data-testid="stSidebar"] .stSelectbox [data-baseweb="select"] *,
section[data-testid="stSidebar"] input {
    color: #111827 !important;
}
span[data-baseweb="tag"] {
    background-color: #D9DDE3 !important;
    color: #111827 !important;
    border-radius: 8px !important;
}
span[data-baseweb="tag"] svg {
    color: #111827 !important;
}
.hero {
    background: #000000;
    color: #ffffff;
    padding: 28px 32px;
    border-radius: 22px;
    margin-bottom: 24px;
    position: relative;
    overflow: hidden;
    border-bottom: 8px solid #009FE3;
}
.hero::after {
    content: "";
    position: absolute;
    width: 430px;
    height: 430px;
    border-radius: 50%;
    background: linear-gradient(135deg, #009FE3, #6F6F6F);
    right: -180px;
    top: -230px;
    opacity: 0.55;
}
.hero-logo {
    position: relative;
    z-index: 2;
    margin-bottom: 18px;
}
.hero-logo img {
    height: 42px;
    width: auto;
    background: #ffffff;
    padding: 8px 12px;
    border-radius: 8px;
}
.hero-title {
    position: relative;
    z-index: 2;
    font-size: 42px;
    line-height: 1.05;
    font-weight: 800;
    letter-spacing: -0.04em;
    max-width: 1100px;
}
.hero-subtitle {
    position: relative;
    z-index: 2;
    margin-top: 12px;
    color: #e6e6e6;
    font-size: 17px;
    max-width: 1150px;
}
.hero-meta {
    position: relative;
    z-index: 2;
    margin-top: 18px;
    font-size: 14px;
    color: #d8d8d8;
}
.section-kicker {
    color: #8D96A0;
    font-size: 13px;
    font-weight: 800;
    letter-spacing: .32em;
    text-transform: uppercase;
    margin: 28px 0 18px;
    display: flex;
    align-items: center;
    gap: 14px;
}
.section-kicker::after {
    content: "";
    height: 1px;
    background: #E2E6EA;
    flex: 1;
}
.brand-strip {
    display: flex;
    align-items: center;
    gap: 36px;
    margin: 8px 0 20px 0;
}
.brand-strip img {
    object-fit: contain;
}
.card {
    background: #ffffff;
    border: 1px solid #e6e9ed;
    border-radius: 16px;
    padding: 20px 22px;
    box-shadow: 0 1px 8px rgba(0,0,0,.035);
}
.benchmark-title {
    font-size: 18px;
    font-weight: 800;
    color: #0A2342;
    margin-bottom: 10px;
}
.benchmark-copy {
    color: #6F7782;
    line-height: 1.45;
    font-size: 14px;
}
.benchmark-metric {
    font-size: 28px;
    color: #0A2342;
    margin-top: 16px;
}
.insight-card {
    background: #ffffff;
    border: 1px solid #e6e9ed;
    border-radius: 14px;
    padding: 20px 22px;
    min-height: 235px;
    box-shadow: 0 1px 8px rgba(0,0,0,.035);
    overflow: hidden;
}
.insight-card.opportunity {
    border-left: 5px solid #12C76B;
}
.insight-card.risk {
    border-left: 5px solid #FFB000;
}
.insight-card.intelligence {
    border-left: 5px solid #2563EB;
}
.market-flag {
    display: inline-block;
    background: #F3F7FF;
    color: #005F9E;
    border: 1px solid #CFE8FF;
    border-radius: 999px;
    padding: 4px 10px;
    font-size: 12px;
    font-weight: 800;
    margin-bottom: 10px;
}
.insight-label {
    font-size: 12px;
    letter-spacing: .22em;
    text-transform: uppercase;
    font-weight: 800;
    margin-bottom: 12px;
}
.insight-label.opportunity {
    color: #12C76B;
}
.insight-label.risk {
    color: #FFB000;
}
.insight-label.intelligence {
    color: #2563EB;
}
.insight-title {
    font-size: 18px;
    line-height: 1.25;
    font-weight: 800;
    color: #0A2342;
    margin-bottom: 10px;
}
.insight-copy {
    color: #8B95A1;
    font-size: 14px;
    line-height: 1.55;
    min-height: 90px;
}
.insight-metric {
    color: #0A2342;
    font-size: 27px;
    margin-top: 14px;
    font-weight: 600;
}
.tag {
    display: inline-block;
    background: #EEF3FF;
    color: #2563EB;
    border-radius: 5px;
    padding: 4px 10px;
    font-size: 12px;
    margin-top: 10px;
}
.footer-note {
    margin-top: 36px;
    padding: 18px 20px;
    background: #F3F3F3;
    border-left: 6px solid #009FE3;
    border-radius: 12px;
    font-size: 13px;
    color: #0A2342;
}
div[data-testid="stMetricValue"] {
    color: #000000;
}
div[data-testid="stPlotlyChart"] {
    background: #ffffff;
    border: 1px solid #e6e9ed;
    border-radius: 14px;
    padding: 10px;
    box-shadow: 0 1px 8px rgba(0,0,0,.035);
}

.tl-detail-card {
    background: #ffffff;
    border: 1px solid #e6e9ed;
    border-radius: 16px;
    padding: 22px;
    box-shadow: 0 1px 8px rgba(0,0,0,.035);
    min-height: 260px;
}
.tl-detail-logo {
    height: 46px;
    display: flex;
    align-items: center;
    margin-bottom: 18px;
}
.tl-detail-logo img {
    max-height: 42px;
    max-width: 230px;
    object-fit: contain;
}
.tl-detail-grid {
    display: grid;
    grid-template-columns: repeat(3, minmax(0, 1fr));
    gap: 14px;
}
.tl-detail-metric {
    background: #F7F9FC;
    border-radius: 12px;
    padding: 14px;
    min-height: 120px;
}
.tl-detail-label {
    color: #6F7782;
    font-size: 13px;
    font-weight: 700;
    margin-bottom: 8px;
}
.tl-detail-value {
    color: #000000;
    font-size: 30px;
    line-height: 1.1;
    font-weight: 800;
    margin-bottom: 8px;
}
.tl-detail-delta-pos {
    display: inline-block;
    background: #DDF8EC;
    color: #128A3A;
    padding: 4px 8px;
    border-radius: 999px;
    font-size: 13px;
    font-weight: 800;
}
.tl-detail-delta-neg {
    display: inline-block;
    background: #FFE5EF;
    color: #D33F49;
    padding: 4px 8px;
    border-radius: 999px;
    font-size: 13px;
    font-weight: 800;
}


.start-intro {
    background:#F7F9FC;
    border-left:6px solid #009FE3;
    border-radius:14px;
    padding:20px 22px;
    color:#0A2342;
    font-size:16px;
    line-height:1.6;
    margin-bottom:10px;
}
.start-shell {
    position: relative;
    min-height: 980px;
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-radius:18px;
    box-shadow: 0 1px 8px rgba(0,0,0,.035);
    overflow:hidden;
    margin-bottom: 20px;
}
.factor-box {
    position:absolute;
    width:240px;
    min-height:96px;
    background:transparent;
    border-left:8px solid #B8EC63;
    padding:18px 16px;
    display:flex;
    align-items:center;
    justify-content:center;
    text-align:center;
    color:#111827;
    font-size:15px;
    line-height:1.3;
    font-weight:600;
    box-shadow:0 1px 4px rgba(0,0,0,.04);
}
.funnel-wrap {
    position:absolute;
    left:50%;
    top:70px;
    transform:translateX(-50%);
    width:390px;
    z-index:2;
}
.funnel-label-top, .funnel-label-bottom {
    text-align:center;
    font-weight:800;
    color:#0A2342;
    margin-bottom:10px;
}
.funnel-seg {
    margin: 0 auto 10px auto;
    color:#ffffff;
    display:flex;
    align-items:center;
    justify-content:center;
    text-align:center;
    font-weight:700;
    border-radius:12px;
    padding:10px 14px;
    background:linear-gradient(135deg, #009FE3, #005F9E);
    box-shadow:0 1px 4px rgba(0,0,0,.08);
}
.funnel-top { width: 360px; height: 62px; background:linear-gradient(135deg, #12C76B, #009FE3); }
.funnel-mid1 { width: 300px; height: 58px; clip-path: polygon(6% 0, 94% 0, 100% 100%, 0 100%); }
.funnel-mid2 { width: 245px; height: 58px; clip-path: polygon(8% 0, 92% 0, 100% 100%, 0 100%); background:linear-gradient(135deg, #007EB5, #0A2342); }
.funnel-mid3 { width: 195px; height: 54px; clip-path: polygon(10% 0, 90% 0, 100% 100%, 0 100%); background:linear-gradient(135deg, #0A2342, #005F9E); }
.funnel-bottom { width: 138px; height: 62px; background:linear-gradient(135deg, #FFB000, #12C76B); }
.funnel-note {
    text-align:center;
    color:#6F7782;
    font-size:13px;
    line-height:1.45;
    margin-top:10px;
}
.method-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-radius:14px;
    padding:18px 20px;
    box-shadow:0 1px 8px rgba(0,0,0,.035);
    height:100%;
}
.method-title {
    font-size:16px;
    font-weight:800;
    color:#0A2342;
    margin-bottom:8px;
}
.method-copy {
    color:#6F7782;
    font-size:14px;
    line-height:1.55;
}
.oem-pill-wrap {
    display:flex;
    flex-wrap:wrap;
    gap:10px;
}
.oem-pill {
    background:#EEF3F8;
    color:#0A2342;
    border:1px solid #D9E2EA;
    padding:8px 12px;
    border-radius:999px;
    font-size:13px;
    font-weight:700;
}


.start-factors-grid {
    display:grid;
    grid-template-columns: repeat(3, minmax(0, 1fr));
    gap:16px;
    margin-bottom:20px;
}
.start-factor-card {
    background:#F5F6F7;
    border-left:8px solid #B8EC63;
    border-radius:8px;
    min-height:72px;
    padding:14px 16px;
    display:flex;
    align-items:center;
    justify-content:center;
    text-align:center;
    color:#111827;
    font-size:18px;
    line-height:1.3;
    font-weight:600;
    box-shadow:0 1px 4px rgba(0,0,0,.04);
}
@media (max-width: 900px) {
    .start-factors-grid { grid-template-columns: 1fr; }
}


.methodology-grid {
    display:grid;
    grid-template-columns: repeat(2, minmax(0, 1fr));
    gap:18px;
    margin-bottom:20px;
}
.methodology-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-radius:16px;
    padding:18px;
    box-shadow:0 1px 8px rgba(0,0,0,.035);
    height:100%;
}
.methodology-card h4 {
    margin:0 0 10px 0;
    font-size:20px;
    color:#0A2342;
}
.methodology-callout {
    display:inline-block;
    background:#EEF3F8;
    color:#0A2342;
    border-left:5px solid #009FE3;
    padding:8px 12px;
    border-radius:8px;
    font-size:13px;
    font-weight:700;
    margin-bottom:14px;
}
.methodology-formula {
    background:#F7F9FC;
    border:1px solid #E3E8EE;
    border-radius:10px;
    padding:12px 14px;
    font-size:14px;
    color:#111827;
    margin:10px 0 12px 0;
    line-height:1.45;
}
.methodology-card ul {
    margin: 10px 0 0 18px;
    color:#6F7782;
    line-height:1.6;
}
@media (max-width: 900px) {
    .methodology-grid { grid-template-columns: 1fr; }
}


.usecase-grid {
    display:grid;
    grid-template-columns: repeat(2, minmax(0, 1fr));
    gap:18px;
    margin-bottom:14px;
}
.usecase-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-radius:16px;
    padding:18px;
    box-shadow:0 1px 8px rgba(0,0,0,.035);
    min-height:255px;
}
.usecase-card.nmsc {
    border-left:7px solid #009FE3;
}
.usecase-card.tme {
    border-left:7px solid #6F6F6F;
}
.usecase-card.shared {
    border-left:7px solid #B8EC63;
}
.usecase-audience {
    display:inline-block;
    background:#EEF3F8;
    color:#0A2342;
    border-radius:999px;
    padding:5px 10px;
    font-size:12px;
    font-weight:800;
    margin-bottom:12px;
}
.usecase-title {
    font-size:20px;
    font-weight:800;
    color:#0A2342;
    margin-bottom:10px;
}
.usecase-copy {
    color:#6F7782;
    font-size:14px;
    line-height:1.55;
    margin-bottom:14px;
}
.usecase-report {
    background:#F7F9FC;
    border-radius:10px;
    padding:12px 14px;
    color:#111827;
    font-size:14px;
    line-height:1.5;
}
.usecase-report b {
    color:#0A2342;
}
.market-link-grid {
    display:grid;
    grid-template-columns: repeat(5, minmax(0, 1fr));
    gap:12px;
    margin:12px 0 22px 0;
}
.market-link-card {
    display:block;
    text-decoration:none !important;
    background:#000000;
    color:#ffffff !important;
    border-radius:14px;
    padding:16px;
    text-align:center;
    font-weight:800;
    box-shadow:0 1px 8px rgba(0,0,0,.08);
}
.market-link-card:hover {
    background:#009FE3;
    color:#000000 !important;
}
.usecase-note {
    background:#F7F9FC;
    border-left:6px solid #009FE3;
    border-radius:14px;
    padding:18px 20px;
    color:#0A2342;
    line-height:1.55;
    margin-bottom:10px;
}
@media (max-width: 900px) {
    .usecase-grid { grid-template-columns: 1fr; }
    .market-link-grid { grid-template-columns: 1fr; }
}


.cluster-legend {
    display:flex;
    flex-wrap:wrap;
    gap:10px;
    margin: 8px 0 18px 0;
}
.cluster-chip {
    display:flex;
    align-items:center;
    gap:8px;
    background:#F7F9FC;
    border:1px solid #E1E7EF;
    border-radius:999px;
    padding:8px 12px;
    color:#0A2342;
    font-size:13px;
}
.cluster-chip span:last-child {
    color:#8B95A1;
    font-weight:500;
}
.cluster-dot {
    display:inline-block;
    width:11px;
    height:11px;
    border-radius:4px;
}


.bubble-control-note {
    background:#F7F9FC;
    border-left:6px solid #009FE3;
    border-radius:14px;
    padding:16px 18px;
    color:#0A2342;
    line-height:1.5;
    margin: 8px 0 18px 0;
}
.bubble-side-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-radius:16px;
    padding:18px;
    box-shadow:0 1px 8px rgba(0,0,0,.035);
    max-height:720px;
    overflow:auto;
}
.bubble-side-title {
    color:#0A2342;
    font-size:16px;
    font-weight:800;
    margin-bottom:12px;
}
.bubble-row {
    display:grid;
    grid-template-columns: 22px 1fr auto;
    gap:8px;
    align-items:center;
    border-bottom:1px solid #EEF2F6;
    padding:9px 0;
}
.bubble-row:last-child {
    border-bottom:none;
}
.bubble-row-dot {
    width:12px;
    height:12px;
    border-radius:4px;
}
.bubble-row-brand {
    color:#0A2342;
    font-weight:700;
    font-size:13px;
}
.bubble-row-meta {
    color:#8B95A1;
    font-size:12px;
}
.bubble-row-value {
    color:#0A2342;
    font-weight:800;
    font-size:13px;
    text-align:right;
}


.assistant-shell {
    background:#000000;
    color:#ffffff;
    border-radius:20px;
    padding:24px;
    border-bottom:7px solid #009FE3;
    margin: 10px 0 22px 0;
}
.assistant-title {
    font-size:26px;
    font-weight:800;
    letter-spacing:-0.02em;
    margin-bottom:8px;
}
.assistant-copy {
    color:#D9DDE3;
    font-size:15px;
    line-height:1.5;
    max-width:1100px;
}
.assistant-hint-row {
    display:flex;
    flex-wrap:wrap;
    gap:8px;
    margin-top:14px;
}
.assistant-hint {
    background:#ffffff;
    color:#0A2342;
    border-radius:999px;
    padding:7px 11px;
    font-size:12px;
    font-weight:800;
}
.assistant-result-grid {
    display:grid;
    grid-template-columns: repeat(3, minmax(0, 1fr));
    gap:16px;
    margin: 16px 0 22px 0;
}
.assistant-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-radius:16px;
    padding:20px;
    box-shadow:0 1px 8px rgba(0,0,0,.035);
    min-height:190px;
}
.assistant-card.primary {
    border-left:7px solid #009FE3;
}
.assistant-card.positive {
    border-left:7px solid #12C76B;
}
.assistant-card.risk {
    border-left:7px solid #FF5C8A;
}
.assistant-card.neutral {
    border-left:7px solid #6F6F6F;
}
.assistant-card-title {
    font-size:17px;
    line-height:1.25;
    color:#0A2342;
    font-weight:800;
    margin-bottom:8px;
}
.assistant-card-copy {
    color:#6F7782;
    font-size:14px;
    line-height:1.55;
    margin-bottom:12px;
}
.assistant-card-metric {
    color:#0A2342;
    font-size:28px;
    line-height:1.1;
    font-weight:800;
    margin-top:8px;
}
.assistant-card-tag {
    display:inline-block;
    background:#EEF3FF;
    color:#2563EB;
    border-radius:6px;
    padding:4px 9px;
    font-size:12px;
    font-weight:700;
    margin-top:10px;
}
@media (max-width: 900px) {
    .assistant-result-grid { grid-template-columns: 1fr; }
}


/* Data Assistant premium concierge treatment */
.da-hero {
    position: relative;
    min-height: 320px;
    border-radius: 26px;
    overflow: hidden;
    margin: 6px 0 28px 0;
    background:
        linear-gradient(90deg, rgba(0,0,0,0.72) 0%, rgba(0,0,0,0.48) 42%, rgba(10,35,66,0.30) 100%),
        radial-gradient(circle at 75% 20%, rgba(255,255,255,0.26), rgba(255,255,255,0) 32%),
        linear-gradient(135deg, #050505 0%, #2A2F36 48%, #C9B7AA 100%);
    border-bottom: 7px solid #009FE3;
    box-shadow: 0 12px 38px rgba(0,0,0,0.16);
}
.da-hero::before {
    content: "";
    position:absolute;
    left:0;
    top:0;
    width:38%;
    height:100%;
    background:
        repeating-linear-gradient(90deg, rgba(255,255,255,0.08) 0px, rgba(255,255,255,0.08) 4px, transparent 4px, transparent 18px);
    opacity: .72;
}
.da-hero-content {
    position: relative;
    z-index: 2;
    padding: 28px 38px 28px 38px;
    text-align: center;
    color:#ffffff;
}
.da-eyebrow {
    display:inline-block;
    background:rgba(255,255,255,0.14);
    border:1px solid rgba(255,255,255,0.24);
    border-radius:999px;
    padding:8px 14px;
    font-size:12px;
    font-weight:800;
    letter-spacing:.14em;
    text-transform:uppercase;
    margin-bottom:14px;
}
.da-title {
    font-size:46px;
    line-height:0.98;
    font-weight:800;
    letter-spacing:-0.055em;
    margin:0 auto 18px auto;
    max-width:900px;
}
.da-subtitle {
    font-size:17px;
    line-height:1.45;
    color:rgba(255,255,255,0.88);
    max-width:1040px;
    margin:0 auto 22px auto;
    font-weight:500;
}
.da-popular-label {
    color:#ffffff;
    font-size:15px;
    font-weight:800;
    letter-spacing:.16em;
    text-transform:uppercase;
    margin-bottom:10px;
}
.da-popular-grid {
    display:flex;
    justify-content:center;
    flex-wrap:wrap;
    gap:18px;
}
.da-popular-pill {
    background:#ffffff;
    color:#101820;
    border-radius:12px;
    padding:12px 18px;
    font-size:16px;
    font-weight:700;
    box-shadow:0 3px 14px rgba(0,0,0,.16);
    min-width:230px;
}
.da-privacy {
    position:absolute;
    bottom:14px;
    left:0;
    right:0;
    text-align:center;
    color:rgba(255,255,255,0.72);
    font-size:13px;
    letter-spacing:.02em;
    z-index:2;
}
.da-search-wrap {
    background:#ffffff;
    border:1.5px solid rgba(167, 100, 70, .55);
    border-radius:999px;
    padding:8px 18px 8px 20px;
    display:flex;
    align-items:center;
    gap:14px;
    min-height:68px;
    box-shadow:0 8px 30px rgba(0,0,0,.08);
    margin: 6px 0 26px 0;
}
.da-search-icon {
    color:#B87554;
    font-size:24px;
    font-weight:900;
}
.da-search-wrap div[data-testid="stTextInput"] {
    flex:1;
}
.da-search-wrap div[data-testid="stTextInput"] > label {
    display:none;
}
.da-search-wrap div[data-testid="stTextInput"] input {
    border:0 !important;
    box-shadow:none !important;
    font-size:21px !important;
    color:#101820 !important;
    background:#ffffff !important;
}
.da-search-wrap div[data-testid="stTextInput"] input:focus {
    border:0 !important;
    box-shadow:none !important;
}
.da-context-bar {
    display:flex;
    justify-content:space-between;
    align-items:center;
    gap:14px;
    margin: 6px 0 18px 0;
    color:#6F7782;
    font-size:14px;
}
.da-context-chip {
    background:#F7F9FC;
    border:1px solid #E6E9ED;
    color:#0A2342;
    border-radius:999px;
    padding:8px 12px;
    font-weight:800;
}
.assistant-result-grid {
    display:grid;
    grid-template-columns: repeat(3, minmax(0, 1fr));
    gap:18px;
    margin: 16px 0 22px 0;
}
.assistant-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-radius:18px;
    padding:18px;
    box-shadow:0 4px 18px rgba(10,35,66,.06);
    min-height:175px;
}
.assistant-card.primary {
    border-left:7px solid #009FE3;
}
.assistant-card.positive {
    border-left:7px solid #12C76B;
}
.assistant-card.risk {
    border-left:7px solid #FF5C8A;
}
.assistant-card.neutral {
    border-left:7px solid #B87554;
}
.assistant-card-title {
    font-size:18px;
    line-height:1.25;
    color:#0A2342;
    font-weight:850;
    margin-bottom:10px;
}
.assistant-card-copy {
    color:#6F7782;
    font-size:14px;
    line-height:1.58;
    margin-bottom:14px;
}
.assistant-card-metric {
    color:#000000;
    font-size:27px;
    line-height:1.1;
    font-weight:850;
    margin-top:8px;
}
.assistant-card-tag {
    display:inline-block;
    background:#F3F6FA;
    color:#0A2342;
    border-radius:8px;
    padding:5px 10px;
    font-size:12px;
    font-weight:800;
    margin-top:12px;
}
@media (max-width: 900px) {
    .da-hero { min-height: 620px; }
    .da-hero-content { padding: 40px 22px; }
    .da-title { font-size:46px; }
    .da-subtitle { font-size:18px; margin-bottom:40px; }
    .da-popular-pill { width:100%; min-width:0; }
    .assistant-result-grid { grid-template-columns: 1fr; }
}




/* Data Assistant input cleanup */
div[data-testid="stTextInput"] input {
    min-height: 48px !important;
    font-size: 18px !important;
    padding: 10px 14px !important;
}


/* Simplified Data Assistant */
.da-simple {
    margin: 8px 0 22px 0;
}
.da-simple-title {
    color:#0A2342;
    font-size:44px;
    line-height:1.05;
    font-weight:850;
    letter-spacing:-0.04em;
    margin-bottom:10px;
}
.da-simple-copy {
    color:#6F7782;
    font-size:18px;
    line-height:1.45;
    max-width:1120px;
    margin-bottom:22px;
}
.da-simple-label {
    color:#8B95A1;
    font-size:13px;
    font-weight:850;
    letter-spacing:.16em;
    text-transform:uppercase;
    margin-bottom:12px;
}
.da-simple-questions {
    display:flex;
    flex-wrap:wrap;
    gap:12px;
    margin-bottom:18px;
}
.da-simple-question {
    background:#ffffff;
    color:#0A2342;
    border:1px solid #E1E7EF;
    border-radius:14px;
    padding:13px 16px;
    font-size:15px;
    font-weight:800;
    box-shadow:0 2px 10px rgba(10,35,66,.05);
}
.da-simple-note {
    color:#8B95A1;
    font-size:13px;
    margin:8px 0 18px 0;
}
@media (max-width: 900px) {
    .da-simple-title { font-size:34px; }
    .da-simple-copy { font-size:16px; }
    .da-simple-question { width:100%; }
}


/* Selected OEM logo sizing cleanup */
.tl-logo-strip img[src*="ToyotaLogo"],
.tl-logo-strip img[src*="BHUB_Logo_ToyotaLogo"],
.brand-strip img[src*="ToyotaLogo"],
.brand-strip img[src*="BHUB_Logo_ToyotaLogo"],
img[src*="BHUB_Logo_ToyotaLogo_01.svg"] {
    background: transparent !important;
    background-color: transparent !important;
    box-shadow: none !important;
    border: 0 !important;
    padding: 0 !important;
    width: 210px !important;
    max-width: 210px !important;
    height: auto !important;
    object-fit: contain !important;
}

.tl-logo-strip,
.brand-strip {
    background: transparent !important;
}


/* Executive Market Summary */
.exec-kpi-grid {
    display:grid;
    grid-template-columns: repeat(4, minmax(0, 1fr));
    gap:16px;
    margin: 12px 0 22px 0;
}
.exec-kpi-card {
    background:#F7F9FC;
    border-radius:18px;
    padding:20px;
    border:1px solid #E6E9ED;
    min-height:145px;
}
.exec-kpi-label {
    color:#6F7782;
    font-size:14px;
    font-weight:800;
    margin-bottom:10px;
}
.exec-kpi-value {
    color:#000000;
    font-size:34px;
    line-height:1.05;
    font-weight:850;
}
.exec-kpi-delta {
    display:inline-block;
    margin-top:12px;
    border-radius:999px;
    padding:5px 10px;
    font-weight:800;
    font-size:13px;
    background:#EEF3F8;
    color:#0A2342;
}
.share-grid {
    display:grid;
    grid-template-columns: repeat(2, minmax(0, 1fr));
    gap:18px;
    margin: 12px 0 24px 0;
}
.share-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-radius:18px;
    padding:22px;
    box-shadow:0 2px 14px rgba(10,35,66,.04);
}
.share-title {
    font-size:20px;
    font-weight:850;
    color:#0A2342;
    margin-bottom:14px;
}
.share-metric-row {
    display:flex;
    justify-content:space-between;
    gap:12px;
    border-top:1px solid #EEF2F6;
    padding:12px 0;
}
.share-metric-row:first-of-type {
    border-top:0;
}
.share-label {
    color:#6F7782;
    font-weight:700;
}
.share-value {
    color:#000000;
    font-weight:850;
}
.takeaway-grid {
    display:grid;
    grid-template-columns: repeat(5, minmax(0, 1fr));
    gap:14px;
    margin: 12px 0 22px 0;
}
.takeaway-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-left:7px solid #009FE3;
    border-radius:16px;
    padding:18px;
    min-height:185px;
    box-shadow:0 2px 14px rgba(10,35,66,.04);
}
.takeaway-card.risk { border-left-color:#FF5C8A; }
.takeaway-card.positive { border-left-color:#12C76B; }
.takeaway-card.neutral { border-left-color:#B87554; }
.takeaway-label {
    font-size:11px;
    letter-spacing:.14em;
    text-transform:uppercase;
    font-weight:850;
    color:#8B95A1;
    margin-bottom:10px;
}
.takeaway-title {
    color:#0A2342;
    font-size:16px;
    font-weight:850;
    line-height:1.25;
    margin-bottom:9px;
}
.takeaway-copy {
    color:#6F7782;
    font-size:13px;
    line-height:1.5;
}
@media (max-width: 1100px) {
    .exec-kpi-grid, .takeaway-grid { grid-template-columns: repeat(2, minmax(0, 1fr)); }
}
@media (max-width: 800px) {
    .exec-kpi-grid, .share-grid, .takeaway-grid { grid-template-columns: 1fr; }
}


/* Legacy logo render fix */
img[src^="data:image/png"][alt="Toyota logo"] {
    background: transparent !important;
    border: 0 !important;
    box-shadow: none !important;
    padding: 0 !important;
    object-fit: contain !important;
}
.brand-strip img[alt="Toyota logo"],
.tl-detail-logo img[src^="data:image/png"] {
    max-height: 64px !important;
    width: auto !important;
}


/* Market Summary colour coding */
.exec-kpi-delta.positive, .share-value.positive {
    background:#DDF8EC !important;
    color:#12C76B !important;
}
.exec-kpi-delta.negative, .share-value.negative {
    background:#FFE5EF !important;
    color:#FF2F6D !important;
}
.exec-kpi-delta.neutral, .share-value.neutral {
    background:#EEF2F6 !important;
    color:#0A2342 !important;
}
.share-value.positive, .share-value.negative, .share-value.neutral {
    display:inline-block;
    border-radius:999px;
    padding:5px 10px;
    font-weight:850;
}
.market-summary-note {
    color:#6F7782;
    font-size:13px;
    margin-top:-6px;
    margin-bottom:12px;
}


/* v61 executive cleanup */
.exec-market-grid {display:grid;grid-template-columns: repeat(5, minmax(0, 1fr));gap:14px;margin:10px 0 24px 0;}
.exec-market-card {background:#ffffff;border:1px solid #E6E9ED;border-radius:16px;padding:16px;box-shadow:0 2px 12px rgba(10,35,66,.04);}
.exec-market-card.overall {background:#F7F9FC;border-left:7px solid #009FE3;}
.exec-market-title {color:#0A2342;font-size:16px;font-weight:850;margin-bottom:10px;}
.exec-market-row {display:flex;justify-content:space-between;gap:12px;border-top:1px solid #EEF2F6;padding:8px 0;font-size:13px;}
.exec-market-row:first-of-type {border-top:0;}
.exec-market-label {color:#6F7782;font-weight:700;}
.exec-market-value {color:#000000;font-weight:850;text-align:right;}
.html-table-wrap {width:100%;overflow-x:auto;border:1px solid #E6E9ED;border-radius:16px;margin:10px 0 22px 0;}
.html-table {width:100%;border-collapse:collapse;font-size:14px;}
.html-table th {background:#F7F9FC;color:#6F7782;text-align:left;padding:12px;border-bottom:1px solid #E6E9ED;font-weight:850;}
.html-table td {padding:12px;border-bottom:1px solid #EEF2F6;color:#303642;vertical-align:middle;}
.html-table tr:last-child td {border-bottom:0;}
.html-table .sticky-brand {position:sticky;left:0;background:#ffffff;z-index:2;font-weight:850;color:#0A2342;}
.html-table th.sticky-brand {background:#F7F9FC;z-index:3;}
.table-highlight td, .table-highlight .sticky-brand {background:#FFF7E8 !important;}
.score-top-grid {display:grid;grid-template-columns: repeat(2, minmax(0, 1fr));gap:16px;margin:12px 0 22px 0;}
.score-podium-card {background:#ffffff;border:1px solid #E6E9ED;border-radius:18px;padding:20px;box-shadow:0 2px 14px rgba(10,35,66,.05);}
.score-podium-title {color:#0A2342;font-weight:850;font-size:20px;margin-bottom:12px;}
.score-podium-row {display:grid;grid-template-columns:36px 1fr auto;gap:10px;align-items:center;padding:10px 0;border-top:1px solid #EEF2F6;}
.score-podium-row:first-of-type {border-top:0;}
.score-podium-brand {font-weight:850;color:#0A2342;}
.score-podium-value {font-weight:850;color:#000000;}
.methodology-source-row {display:flex;gap:10px;flex-wrap:wrap;margin-top:12px;}
.source-logo-pill {display:inline-block;background:#ffffff;border:1px solid #E1E7EF;border-radius:10px;padding:8px 12px;color:#0A2342;font-weight:850;box-shadow:0 1px 6px rgba(10,35,66,.04);}
.flow-usecase-grid {display:grid;grid-template-columns: repeat(3, minmax(0, 1fr));gap:16px;margin:16px 0 24px 0;}
.flow-usecase-card {background:#ffffff;border:1px solid #E6E9ED;border-radius:18px;padding:20px;box-shadow:0 2px 14px rgba(10,35,66,.05);min-height:185px;}
.flow-number {width:34px;height:34px;border-radius:50%;background:#000;color:#fff;display:flex;align-items:center;justify-content:center;font-weight:850;margin-bottom:12px;}
.flow-title {color:#0A2342;font-weight:850;font-size:18px;line-height:1.25;}
.flow-copy {color:#6F7782;font-size:14px;line-height:1.5;margin-top:10px;}
.bubble-key {background:#F7F9FC;border:1px solid #E6E9ED;border-radius:14px;padding:12px 14px;color:#0A2342;font-weight:750;margin:8px 0 14px 0;}
@media (max-width:1100px){.exec-market-grid,.flow-usecase-grid{grid-template-columns:repeat(2,minmax(0,1fr));}.score-top-grid{grid-template-columns:1fr;}}
@media (max-width:800px){.exec-market-grid,.flow-usecase-grid{grid-template-columns:1fr;}}


/* v62 refinements */
.market-card-full { margin: 10px 0 16px 0; }
.market-five-grid {
    display:grid;
    grid-template-columns: repeat(5, minmax(0, 1fr));
    gap:14px;
    margin: 10px 0 24px 0;
}
.flag {
    display:inline-block;
    margin-right:7px;
}
.bubble-page-grid {
    display:grid;
    grid-template-columns: minmax(0, 4.8fr) 320px;
    gap:18px;
    align-items:start;
}
.bubble-filter-panel {
    position:sticky;
    top:18px;
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-radius:18px;
    padding:18px;
    box-shadow:0 2px 14px rgba(10,35,66,.05);
}
.bubble-filter-title {
    color:#0A2342;
    font-size:18px;
    font-weight:850;
    margin-bottom:12px;
}
.usecase-toggle-note {
    background:#F7F9FC;
    border:1px solid #E6E9ED;
    border-radius:16px;
    padding:16px 18px;
    margin:10px 0 18px 0;
    color:#6F7782;
}
.dashboard-link-grid {
    display:grid;
    grid-template-columns: repeat(5, minmax(0, 1fr));
    gap:12px;
    margin: 12px 0 24px 0;
}
.dashboard-link-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-radius:16px;
    padding:16px;
    box-shadow:0 2px 12px rgba(10,35,66,.04);
    color:#0A2342 !important;
    text-decoration:none !important;
    font-weight:850;
}
.dashboard-link-card span {
    display:block;
    color:#8B95A1;
    font-weight:700;
    font-size:12px;
    margin-top:4px;
}
.brand-market-grid {
    display:grid;
    grid-template-columns: repeat(5, minmax(0, 1fr));
    gap:14px;
    margin: 12px 0 24px 0;
}
.brand-market-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-left:7px solid #009FE3;
    border-radius:16px;
    padding:16px;
    box-shadow:0 2px 12px rgba(10,35,66,.04);
}
.brand-market-title { color:#0A2342; font-weight:850; font-size:16px; margin-bottom:10px; }
.brand-market-value { color:#000000; font-size:26px; font-weight:850; line-height:1.05; margin-bottom:8px; }
.brand-market-meta { color:#6F7782; font-size:13px; line-height:1.45; }
@media (max-width: 1200px) {
    .market-five-grid, .brand-market-grid, .dashboard-link-grid { grid-template-columns: repeat(2, minmax(0, 1fr)); }
    .bubble-page-grid { grid-template-columns: 1fr; }
    .bubble-filter-panel { position:relative; top:0; }
}
@media (max-width: 800px) {
    .market-five-grid, .brand-market-grid, .dashboard-link-grid { grid-template-columns: 1fr; }
}


/* v63 refinements */
.mm5-kpi-grid {
    display:grid;
    grid-template-columns: repeat(5, minmax(0, 1fr));
    gap:14px;
    margin: 10px 0 18px 0;
}
.mm5-kpi-card {
    background:#F7F9FC;
    border:1px solid #E6E9ED;
    border-radius:18px;
    padding:18px;
    box-shadow:0 2px 12px rgba(10,35,66,.04);
}
.mm5-kpi-card.primary { border-left:7px solid #009FE3; }
.mm5-kpi-label { color:#6F7782; font-weight:850; font-size:13px; margin-bottom:9px; }
.mm5-kpi-value { color:#000; font-size:30px; font-weight:850; line-height:1.05; }
.brand-logo-header {
    display:flex;
    align-items:center;
    gap:22px;
    margin: 8px 0 16px 0;
    min-height:74px;
}
.brand-logo-header img {
    max-height:66px;
    width:auto;
    object-fit:contain;
}
.brand-insight-card-grid {
    display:grid;
    grid-template-columns: repeat(5, minmax(0, 1fr));
    gap:14px;
    margin: 12px 0 24px 0;
}
.brand-insight-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-left:7px solid #009FE3;
    border-radius:18px;
    padding:16px;
    box-shadow:0 2px 12px rgba(10,35,66,.04);
    min-height:205px;
}
.brand-insight-card.competitor { border-left-color:#B87554; background:#FBFCFD; }
.brand-insight-title { color:#0A2342; font-weight:850; font-size:16px; margin-bottom:12px; }
.brand-insight-label { color:#6F7782; font-size:12px; font-weight:800; text-transform:uppercase; letter-spacing:.08em; margin-top:10px; }
.brand-insight-main { color:#000; font-size:30px; font-weight:850; line-height:1.05; margin:4px 0 8px 0; }
.brand-insight-meta { color:#6F7782; font-size:13px; line-height:1.45; }
.bubble-flex {
    display:grid;
    grid-template-columns: minmax(0, 3fr) minmax(250px, 1fr);
    gap:18px;
    align-items:start;
}
.bubble-controls-card {
    background:#fff;
    border:1px solid #E6E9ED;
    border-radius:18px;
    padding:18px;
    box-shadow:0 2px 14px rgba(10,35,66,.05);
    position:sticky;
    top:18px;
}
.lock-note {
    background:#FFF7E8;
    border:1px solid #F4D8A8;
    color:#0A2342;
    border-radius:12px;
    padding:10px 12px;
    font-weight:750;
    margin:8px 0 14px 0;
}
@media (max-width: 1200px) {
    .mm5-kpi-grid, .brand-insight-card-grid { grid-template-columns: repeat(2, minmax(0, 1fr)); }
    .bubble-flex { grid-template-columns: 1fr; }
    .bubble-controls-card { position:relative; top:0; }
}
@media (max-width: 800px) {
    .mm5-kpi-grid, .brand-insight-card-grid { grid-template-columns: 1fr; }
}


/* v64 bubble controls and scorecard polish */
.bubble-controls-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-radius:18px;
    padding:18px;
    box-shadow:0 2px 14px rgba(10,35,66,.05);
}
.bubble-controls-card .control-title {
    color:#0A2342;
    font-weight:850;
    font-size:18px;
    margin-bottom:12px;
}


.brand-insight-card .brand-insight-meta {
    margin-top:10px;
}
.brand-insight-card .brand-insight-meta span[style*="border-radius"] {
    display:inline-block;
    margin:4px 4px 4px 0;
    padding:6px 10px !important;
}
.brand-insight-main {
    margin-bottom:12px !important;
}


/* v68 Insight Report export */
.report-builder-grid {
    display:grid;
    grid-template-columns: 1.25fr 1fr;
    gap:18px;
    margin: 12px 0 24px 0;
}
.report-builder-card {
    background:#ffffff;
    border:1px solid #E6E9ED;
    border-radius:18px;
    padding:20px;
    box-shadow:0 2px 14px rgba(10,35,66,.05);
}
.report-builder-title {
    color:#0A2342;
    font-size:22px;
    font-weight:850;
    margin-bottom:10px;
}
.report-builder-copy {
    color:#6F7782;
    font-size:15px;
    line-height:1.55;
}
.report-type-grid {
    display:grid;
    grid-template-columns: repeat(2, minmax(0, 1fr));
    gap:16px;
    margin: 14px 0 8px 0;
}
.report-type-card {
    border:1px solid #E1E7EF;
    border-radius:16px;
    padding:18px;
    background:#F7F9FC;
    min-height:150px;
}
.report-type-card strong {
    display:block;
    color:#0A2342;
    font-size:20px;
    margin-bottom:8px;
}
.report-type-card span {
    color:#6F7782;
    font-size:14px;
    line-height:1.45;
}
.report-output-card {
    background:#F7F9FC;
    border:1px solid #E6E9ED;
    border-radius:18px;
    padding:18px;
    margin: 12px 0 18px 0;
}
@media (max-width: 1000px) {
    .report-builder-grid, .report-type-grid { grid-template-columns: 1fr; }
}



/* Valtech brand parity override */
:root {
    --valtech-black: #000000;
    --valtech-white: #ffffff;
    --valtech-offwhite: #F6F5F1;
    --valtech-blue: #003CB3;
    --valtech-red: #FF4B55;
    --valtech-grey: #6F6F6F;
    --valtech-line: #D9D7D1;
}
html, body, [class*="css"], .stApp, button, input, textarea, select {
    font-family: "Valtech", "Helvetica Neue", Helvetica, Arial, sans-serif !important;
    letter-spacing: -0.015em;
}
.stApp {
    background: var(--valtech-white) !important;
    color: var(--valtech-black) !important;
}
.block-container {
    padding-top: 0.75rem !important;
    max-width: 1480px !important;
}
.valtech-topbar {
    display:flex;
    align-items:center;
    justify-content:space-between;
    gap:16px;
    padding: 8px 0 10px 0;
    margin-bottom: 10px;
    border-bottom: 1px solid var(--valtech-line);
    background: var(--valtech-white);
}
.valtech-wordmark {
    color: var(--valtech-black);
    font-size: 34px;
    line-height: 1;
    font-weight: 400;
    letter-spacing: -0.055em;
}
.valtech-star {
    display:inline-block;
    margin-left:6px;
    font-size:30px;
    line-height:1;
    transform: translateY(2px);
}
.valtech-nav {
    display:none;
    align-items:center;
    gap:72px;
    color: var(--valtech-black);
    font-size:24px;
    font-weight:700;
}
.valtech-nav span {
    position:relative;
    padding-top: 10px;
}
.valtech-nav span.active::before {
    content:"";
    position:absolute;
    left:50%;
    top:0;
    transform:translateX(-50%);
    width:24px;
    height:6px;
    background: var(--valtech-red);
}
.hero {
    background: var(--valtech-white) !important;
    color: var(--valtech-black) !important;
    border-radius: 0 !important;
    border: 1px solid var(--valtech-line) !important;
    border-bottom: 6px solid var(--valtech-blue) !important;
    padding: 22px 28px 24px 28px !important;
    margin: 0 0 22px 0 !important;
    box-shadow: none !important;
}
.hero::after { display:none !important; }
.hero-accent-red {
    width: 24px;
    height: 5px;
    background: var(--valtech-red);
    margin-bottom: 12px;
}
.hero-kicker {
    color: var(--valtech-grey);
    font-size: 10px;
    font-weight: 800;
    letter-spacing: .22em;
    text-transform: uppercase;
    margin-bottom: 8px;
}
.hero-title {
    color: var(--valtech-black) !important;
    font-size: clamp(28px, 3.8vw, 54px) !important;
    line-height: 1.02 !important;
    font-weight: 400 !important;
    letter-spacing: -0.065em !important;
    max-width: 900px !important;
}
.hero-subtitle {
    color: var(--valtech-black) !important;
    max-width: 980px !important;
    font-size: 14px !important;
    line-height: 1.35 !important;
    margin-top: 14px !important;
    font-weight: 400 !important;
}
.hero-meta {
    color: var(--valtech-grey) !important;
    margin-top: 14px !important;
    padding-top: 10px;
    border-top: 1px solid var(--valtech-line);
    font-size: 11px !important;
}
.section-kicker {
    color: var(--valtech-black) !important;
    font-size: 16px !important;
    font-weight: 700 !important;
    letter-spacing: .02em !important;
    text-transform: none !important;
    margin: 42px 0 22px !important;
}
.section-kicker::after {
    background: var(--valtech-black) !important;
    opacity: .25;
}
.card,
.method-card,
.methodology-card,
.usecase-card,
.flow-usecase-card,
.report-builder-card,
.report-output-card,
.assistant-card,
.share-card,
.exec-kpi-card,
.exec-market-card,
.score-podium-card,
.bubble-controls-card,
.bubble-filter-panel,
.bubble-side-card,
.brand-insight-card,
.mm5-kpi-card,
.takeaway-card,
.tl-detail-card {
    border: 1px solid var(--valtech-black) !important;
    border-radius: 0 !important;
    box-shadow: none !important;
    background: var(--valtech-white) !important;
}
.start-intro,
.usecase-note,
.usecase-toggle-note,
.bubble-key,
.bubble-control-note,
.footer-note,
.lock-note {
    background: var(--valtech-offwhite) !important;
    border: 1px solid var(--valtech-black) !important;
    border-left: 8px solid var(--valtech-blue) !important;
    border-radius: 0 !important;
    color: var(--valtech-black) !important;
    box-shadow: none !important;
}
.methodology-callout,
.da-context-chip,
.assistant-card-tag,
.tag,
.oem-pill,
.cluster-chip,
.source-logo-pill {
    background: var(--valtech-white) !important;
    color: var(--valtech-black) !important;
    border: 1px solid var(--valtech-black) !important;
    border-radius: 0 !important;
}
.insight-card,
.assistant-card.primary,
.takeaway-card,
.brand-market-card,
.brand-insight-card {
    border-left: 8px solid var(--valtech-blue) !important;
}
.insight-card.risk,
.assistant-card.risk,
.takeaway-card.risk {
    border-left-color: var(--valtech-red) !important;
}
.insight-card.opportunity,
.assistant-card.positive,
.takeaway-card.positive {
    border-left-color: var(--valtech-blue) !important;
}
.report-type-card,
.tl-detail-metric,
.start-factor-card,
.exec-kpi-card.primary,
.mm5-kpi-card.primary {
    background: var(--valtech-offwhite) !important;
    border-radius: 0 !important;
}
.start-factor-card {
    border: 1px solid var(--valtech-black) !important;
    border-left: 8px solid var(--valtech-blue) !important;
}
.benchmark-title,
.insight-title,
.method-title,
.methodology-card h4,
.usecase-title,
.flow-title,
.report-builder-title,
.assistant-card-title,
.share-title,
.takeaway-title,
.exec-market-title,
.score-podium-title,
.brand-insight-title,
.da-simple-title,
.bubble-side-title,
.bubble-filter-title,
.brand-market-title {
    color: var(--valtech-black) !important;
    font-weight: 700 !important;
}
.benchmark-copy,
.insight-copy,
.method-copy,
.usecase-copy,
.flow-copy,
.report-builder-copy,
.assistant-card-copy,
.takeaway-copy,
.brand-market-meta,
.brand-insight-meta,
.market-summary-note,
.da-simple-copy {
    color: var(--valtech-grey) !important;
}
.benchmark-metric,
.insight-metric,
.assistant-card-metric,
.exec-kpi-value,
.mm5-kpi-value,
.brand-insight-main,
.brand-market-value,
.tl-detail-value {
    color: var(--valtech-black) !important;
    font-weight: 400 !important;
    letter-spacing: -0.055em !important;
}
div[data-testid="stPlotlyChart"] {
    border: 1px solid var(--valtech-black) !important;
    border-radius: 0 !important;
    box-shadow: none !important;
    padding: 12px !important;
}
.html-table-wrap {
    border: 1px solid var(--valtech-black) !important;
    border-radius: 0 !important;
}
.html-table th {
    background: var(--valtech-offwhite) !important;
    color: var(--valtech-black) !important;
    border-bottom: 1px solid var(--valtech-black) !important;
}
.html-table td {
    color: var(--valtech-black) !important;
    border-bottom: 1px solid var(--valtech-line) !important;
}
section[data-testid="stSidebar"] {
    background: var(--valtech-black) !important;
    border-right: 8px solid var(--valtech-blue) !important;
}
section[data-testid="stSidebar"] h1,
section[data-testid="stSidebar"] h2,
section[data-testid="stSidebar"] h3,
section[data-testid="stSidebar"] label,
section[data-testid="stSidebar"] p,
section[data-testid="stSidebar"] span {
    color: var(--valtech-white) !important;
}
section[data-testid="stSidebar"] [data-baseweb="radio"] label,
section[data-testid="stSidebar"] [data-baseweb="checkbox"] label {
    color: var(--valtech-white) !important;
}
.stButton > button,
div[data-testid="stDownloadButton"] button {
    background: var(--valtech-black) !important;
    color: var(--valtech-white) !important;
    border: 1px solid var(--valtech-black) !important;
    border-radius: 0 !important;
    font-weight: 700 !important;
}
.stButton > button:hover,
div[data-testid="stDownloadButton"] button:hover {
    background: var(--valtech-blue) !important;
    color: var(--valtech-white) !important;
    border-color: var(--valtech-blue) !important;
}
span[data-baseweb="tag"] {
    background-color: var(--valtech-white) !important;
    color: var(--valtech-black) !important;
    border: 1px solid var(--valtech-black) !important;
    border-radius: 0 !important;
}
[data-baseweb="select"] > div,
div[data-testid="stTextInput"] input,
div[data-testid="stNumberInput"] input {
    border-radius: 0 !important;
    border-color: var(--valtech-black) !important;
}
@media (max-width: 900px) {
    .valtech-topbar { align-items:flex-start; flex-direction:column; padding: 6px 0 8px 0; margin-bottom: 8px; }
    .valtech-wordmark { font-size: 28px !important; }
    .valtech-star { font-size: 25px !important; }
    .valtech-nav { display:none !important; }
    .hero { padding: 16px 18px 18px 18px !important; margin-bottom: 16px !important; }
    .hero-title { font-size: 32px !important; letter-spacing: -0.055em !important; }
    .hero-subtitle { font-size: 13px !important; margin-top: 10px !important; }
    .hero-meta { font-size: 10px !important; }
}
</style>
""",
    unsafe_allow_html=True,
)


# =========================
# Data layer
# =========================

def find_column(columns, candidates):
    lower = {str(c).strip().lower(): c for c in columns}
    for candidate in candidates:
        for key, original in lower.items():
            if candidate in key:
                return original
    raise ValueError(f"Could not find a column matching: {candidates}")


@st.cache_data
def normalise_period(value):
    text_value = str(value).strip()
    if text_value.endswith(".0"):
        text_value = text_value[:-2]
    return text_value


@st.cache_data
def load_data():
    raw = pd.read_excel(DATA_FILE, sheet_name="CORE DATA CUT")
    raw.columns = [str(c).strip() for c in raw.columns]

    brand_col = find_column(raw.columns, ["brand", "oem"])
    market_col = find_column(raw.columns, ["market", "country"])
    period_col = find_column(raw.columns, ["period", "year"])
    sales_col = find_column(raw.columns, ["sales"])
    uv_col = find_column(raw.columns, ["unique visitors (dedupe)", "unique visitors", "unique", "uv"])

    df = raw[[brand_col, market_col, period_col, sales_col, uv_col]].copy()
    df = df.rename(
        columns={
            brand_col: "OEM",
            market_col: "Market",
            period_col: "Period",
            sales_col: "Sales",
            uv_col: "UniqueVisitors",
        }
    )

    df["OEM"] = df["OEM"].astype(str).str.strip()
    df["Market"] = df["Market"].astype(str).str.strip()
    df["Period"] = df["Period"].map(normalise_period)
    df["Sales"] = pd.to_numeric(df["Sales"], errors="coerce")
    df["UniqueVisitors"] = pd.to_numeric(df["UniqueVisitors"], errors="coerce")

    df = df.dropna(subset=["OEM", "Market", "Period", "Sales", "UniqueVisitors"])
    df = df[(df["UniqueVisitors"] > 0) & (df["Sales"] >= 0)]

    market_data = (
        df.groupby(["OEM", "Market", "Period"], as_index=False)
        .agg({"Sales": "sum", "UniqueVisitors": "sum"})
    )
    market_data["ConversionPct"] = market_data["Sales"] / market_data["UniqueVisitors"] * 100

    mm5 = (
        market_data.groupby(["OEM", "Period"], as_index=False)
        .agg({"Sales": "sum", "UniqueVisitors": "sum"})
    )
    mm5["Market"] = "MM5"
    mm5["ConversionPct"] = mm5["Sales"] / mm5["UniqueVisitors"] * 100

    combined = pd.concat(
        [
            mm5[["OEM", "Market", "Period", "Sales", "UniqueVisitors", "ConversionPct"]],
            market_data[["OEM", "Market", "Period", "Sales", "UniqueVisitors", "ConversionPct"]],
        ],
        ignore_index=True,
    )

    return combined.sort_values(["Market", "OEM", "Period"]).reset_index(drop=True)


def apply_comparison_view(data, comparison_name):
    config = COMPARISON_OPTIONS[comparison_name]
    previous_period = config["previous_period"]
    current_period = config["current_period"]

    active = data[data["Period"].isin([previous_period, current_period])].copy()
    active["Year"] = active["Period"].map({previous_period: 2024, current_period: 2025})
    active["PeriodLabel"] = active["Period"].map(
        {
            previous_period: config["previous_label"],
            current_period: config["current_label"],
        }
    )

    return active




def market_year(data, market, year, oems=None):
    df = data[(data["Market"] == market) & (data["Year"] == year)].copy()
    if oems:
        df = df[df["OEM"].isin(oems)]
    return df


def yoy_table(data, market, oems=None):
    d24 = market_year(data, market, 2024, oems).copy()
    d25 = market_year(data, market, 2025, oems).copy()

    if d24.empty or d25.empty:
        return pd.DataFrame()

    d24 = d24.rename(
        columns={
            "Sales": "Sales_2024",
            "UniqueVisitors": "UniqueVisitors_2024",
            "ConversionPct": "ConversionPct_2024",
            "Year": "Year_2024",
        }
    )
    d25 = d25.rename(
        columns={
            "Sales": "Sales_2025",
            "UniqueVisitors": "UniqueVisitors_2025",
            "ConversionPct": "ConversionPct_2025",
            "Year": "Year_2025",
        }
    )

    d24 = d24[["OEM", "Market", "Sales_2024", "UniqueVisitors_2024", "ConversionPct_2024"]]
    d25 = d25[["OEM", "Market", "Sales_2025", "UniqueVisitors_2025", "ConversionPct_2025"]]

    merged = d25.merge(d24, on=["OEM", "Market"], how="inner")
    if merged.empty:
        return pd.DataFrame()

    merged["Sales YoY %"] = (merged["Sales_2025"] / merged["Sales_2024"] - 1) * 100
    merged["Visitors YoY %"] = (merged["UniqueVisitors_2025"] / merged["UniqueVisitors_2024"] - 1) * 100
    merged["Conv Var pp"] = merged["ConversionPct_2025"] - merged["ConversionPct_2024"]
    merged["Visits to Sale 2025"] = merged["UniqueVisitors_2025"] / merged["Sales_2025"]
    merged["Visits to Sale 2024"] = merged["UniqueVisitors_2024"] / merged["Sales_2024"]
    merged["Visits to Sale Var"] = merged["Visits to Sale 2025"] - merged["Visits to Sale 2024"]
    merged = merged.replace([float("inf"), float("-inf")], pd.NA)
    return merged


def selected_or_all(selected, all_oems):
    return selected if selected else all_oems


# Backward-compatible aliases in case any code path uses older names.
get_yoy_table = yoy_table
get_market_data = market_year


# =========================
# Formatting helpers
# =========================

def fmt_int(v):
    if pd.isna(v):
        return "n/a"
    return f"{float(v):,.0f}"


def fmt_pct(v, decimals=1):
    if pd.isna(v):
        return "n/a"
    sign = "+" if float(v) > 0 else ""
    return f"{sign}{float(v):.{decimals}f}%"


def fmt_pp(v):
    if pd.isna(v):
        return "n/a"
    sign = "+" if float(v) > 0 else ""
    return f"{sign}{float(v):.2f}pp"


def fmt_short(v):
    v = float(v)
    if abs(v) >= 1_000_000:
        return f"{v/1_000_000:.2f}M"
    if abs(v) >= 1_000:
        return f"{v/1_000:.1f}K"
    return f"{v:,.0f}"


def fmt_metric_number(v):
    """Format KPI numbers so they do not wrap awkwardly in metric cards."""
    if pd.isna(v):
        return "n/a"
    v = float(v)
    if abs(v) >= 1_000_000:
        return f"{v/1_000_000:.1f}M"
    if abs(v) >= 1_000:
        return f"{v:,.0f}"
    return f"{v:.0f}"


def badge_style(value):
    text = str(value)
    if text.startswith("+"):
        return "background-color:#DDF8EC;color:#12C76B;font-weight:700;border-radius:7px;"
    if text.startswith("-"):
        return "background-color:#FFE5EF;color:#FF2F6D;font-weight:700;border-radius:7px;"
    return "background-color:#EEF2F6;color:#6F6F6F;font-weight:700;border-radius:7px;"


def badge_html(value, suffix="%"):
    try:
        v = float(value)
    except Exception:
        return '<span style="background:#EEF2F6;color:#6F6F6F;padding:5px 9px;border-radius:6px;font-weight:700;">n/a</span>'
    cls_bg = "#DDF8EC" if v > 0 else "#FFE5EF" if v < 0 else "#EEF2F6"
    cls_col = "#12C76B" if v > 0 else "#FF2F6D" if v < 0 else "#6F6F6F"
    sign = "+" if v > 0 else ""
    text = f"{sign}{v:.2f}pp" if suffix == "pp" else f"{sign}{v:.1f}%"
    return f'<span style="background:{cls_bg};color:{cls_col};padding:5px 9px;border-radius:6px;font-weight:700;">{text}</span>'


def get_row(yoy, brand):
    if yoy.empty:
        return None
    row = yoy[yoy["OEM"].str.lower() == brand.lower()]
    return None if row.empty else row.iloc[0]


def preset_selection(preset_name, all_oems):
    if preset_name == "Toyota volume competitors":
        preset = TOYOTA_SET
    elif preset_name == "Lexus premium competitors":
        preset = LEXUS_SET
    else:
        preset = all_oems
    return [x for x in preset if x in all_oems]


def normalise_cluster_oem(oem):
    if oem == "VW":
        return "Volkswagen"
    if oem == "Citroen":
        return "Citroën"
    if oem == "BYD":
        return "BYD Auto"
    return oem


def cluster_for_oem(oem):
    """Return the six-category cluster for an OEM, supporting both dict orientations."""
    if oem in OEM_CLUSTERS and isinstance(OEM_CLUSTERS.get(oem), str):
        return OEM_CLUSTERS.get(oem)
    for cluster, oems in OEM_CLUSTERS.items():
        if isinstance(oems, (list, tuple, set)) and oem in oems:
            return cluster
    return "Other"




def available_oem_categories():
    """Return the six OEM category names used in the dashboard."""
    return [
        "Volume Leaders",
        "EV Challengers",
        "European Mass Market",
        "Luxury & Performance",
        "Asian Mainstream",
        "Chinese New Entrants",
    ]


def oems_for_categories(data, categories):
    """Return OEMs that belong to any selected category."""
    if not categories:
        return sorted(data["OEM"].unique())
    selected = set(categories)
    return sorted([
        oem for oem in data["OEM"].unique()
        if cluster_for_oem(oem) in selected
    ])




def oems_for_clusters(clusters, available_oems):
    if not clusters or "All categories" in clusters:
        return available_oems

    allowed = set()
    for cluster in clusters:
        for brand in OEM_CLUSTERS.get(cluster, []):
            normalised = normalise_cluster_oem(brand)
            for available in available_oems:
                if normalise_cluster_oem(available) == normalised:
                    allowed.add(available)

    return sorted(allowed)


def filter_oems_by_cluster(selected_oems, clusters, available_oems):
    allowed = oems_for_clusters(clusters, available_oems)
    if not selected_oems:
        return allowed
    return sorted([oem for oem in selected_oems if oem in allowed])


def render_cluster_legend(selected_clusters):
    active = [c for c in selected_clusters if c != "All categories"]
    if not active:
        active = list(OEM_CLUSTERS.keys())

    chips = ""
    for cluster in active:
        color = CLUSTER_COLORS.get(cluster, "#6F6F6F")
        desc = CLUSTER_DESCRIPTIONS.get(cluster, "")
        chips += (
            f"<div class='cluster-chip'>"
            f"<span class='cluster-dot' style='background:{color};'></span>"
            f"<b>{cluster}</b><span>{desc}</span>"
            f"</div>"
        )

    st.markdown(f"<div class='cluster-legend'>{chips}</div>", unsafe_allow_html=True)


# =========================
# Visual components
# =========================

def render_hero():
    st.markdown(
        f"""
        <div class="valtech-topbar">
            <div class="valtech-wordmark">Valtech <span class="valtech-star">✳</span></div>
        </div>
        <div class="hero">
            <div class="hero-accent-red"></div>
            <div class="hero-kicker">OEM intelligence</div>
            <div class="hero-title">OEM Macro Conversion Intelligence Report</div>
            <div class="hero-subtitle">
                Website-to-contract conversion intelligence across MM5. Explore how unique visitor demand converts into passenger car sales, compare OEMs consistently by market, and identify cross-market performance opportunities.
            </div>
            <div class="hero-meta">
                <b>Coverage:</b> 2024, 2025 &amp; 2026 (Jan to April) &nbsp;|&nbsp;
                <b>Sales data:</b> www.marklines.com &nbsp;|&nbsp;
                <b>Unique visitors:</b> www.similarweb.com
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def section(title):
    st.markdown(f'<div class="section-kicker">{title}</div>', unsafe_allow_html=True)


def render_footer():
    st.markdown(
        f"""
        <div class="footer-note">
        <b>Definitions and caveats:</b> This is not a total market-share dashboard. It uses passenger car sales from Marklines and <b>monthly deduplicated unique visitor</b> website data from Similarweb. It does not include fleet, LCV or tactical registrations. <b>Website-to-Contract Conversion Rate</b> = passenger car sales / monthly deduplicated unique website visitors. Active comparison: <b>{PREVIOUS_LABEL} vs {CURRENT_LABEL}</b>.
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_brand_strip():
    st.markdown(
        f"""
        <div class="brand-strip">
            <img src="{TOYOTA_LOGO}" style="height:58px; width:auto; object-fit:contain;" alt="Toyota logo">
            <img src="{LEXUS_LOGO}" style="height:30px;" alt="Lexus logo">
        </div>
        """,
        unsafe_allow_html=True,
    )


def delta_class_for_html(value):
    try:
        return "tl-detail-delta-pos" if float(value) >= 0 else "tl-detail-delta-neg"
    except Exception:
        return "tl-detail-delta-pos"


def brand_detail_html(brand, row, logo, logo_height="38px"):
    if row is None:
        return (
            "<div class='tl-detail-card'>"
            f"<div class='tl-detail-logo'><img src='{logo}' style='height:{logo_height};'></div>"
            f"<p>No {brand} data available for this market.</p>"
            "</div>"
        )

    conv_delta_class = delta_class_for_html(row["Conv Var pp"])
    sales_delta_class = delta_class_for_html(row["Sales YoY %"])
    visitor_delta_class = delta_class_for_html(row["Visitors YoY %"])

    return (
        "<div class='tl-detail-card'>"
        f"<div class='tl-detail-logo'><img src='{logo}' style='height:{logo_height};' alt='{brand} logo'></div>"
        "<div class='tl-detail-grid'>"
        "<div class='tl-detail-metric'>"
        f"<div class='tl-detail-label'>Website-to-Contract Conversion Rate</div>"
        f"<div class='tl-detail-value'>{row['ConversionPct_2025']:.2f}%</div>"
        f"<div class='{conv_delta_class}'>{fmt_pp(row['Conv Var pp'])} vs previous period</div>"
        "</div>"
        "<div class='tl-detail-metric'>"
        f"<div class='tl-detail-label'>Passenger sales</div>"
        f"<div class='tl-detail-value'>{fmt_metric_number(row['Sales_2025'])}</div>"
        f"<div class='{sales_delta_class}'>{fmt_pct(row['Sales YoY %'])} vs previous period</div>"
        "</div>"
        "<div class='tl-detail-metric'>"
        f"<div class='tl-detail-label'>Unique visitors</div>"
        f"<div class='tl-detail-value'>{fmt_metric_number(row['UniqueVisitors_2025'])}</div>"
        f"<div class='{visitor_delta_class}'>{fmt_pct(row['Visitors YoY %'])} vs previous period</div>"
        "</div>"
        "</div>"
        "</div>"
    )


def render_brand_detail(data, market):
    section("2025 selected OEM performance vs previous period")
    yoy = yoy_table(data, market, None)
    toyota = get_row(yoy, "Toyota")
    lexus = get_row(yoy, "Lexus")

    cols = st.columns(2)
    with cols[0]:
        st.markdown(brand_detail_html("Toyota", toyota, TOYOTA_LOGO, "42px"), unsafe_allow_html=True)
    with cols[1]:
        st.markdown(brand_detail_html("Lexus", lexus, LEXUS_LOGO, "30px"), unsafe_allow_html=True)


def benchmark_card(title, copy, metric):
    html = (
        "<div class='card'>"
        f"<div class='benchmark-title'>{title}</div>"
        f"<div class='benchmark-copy'>{copy}</div>"
        f"<div class='benchmark-metric'>{metric}</div>"
        "</div>"
    )
    st.markdown(html, unsafe_allow_html=True)


def render_benchmark_cards(data, market):
    section("Selected OEM Website-to-Contract benchmark callouts")
    yoy = yoy_table(data, market, None)
    if yoy.empty:
        st.info("No benchmark data available.")
        return

    cols = st.columns(2)
    for col, brand, cohort in [(cols[0], "Toyota", TOYOTA_SET), (cols[1], "Lexus", LEXUS_SET)]:
        with col:
            row = get_row(yoy, brand)
            cohort_df = yoy[yoy["OEM"].isin(cohort)].copy()
            if row is None or cohort_df.empty:
                st.info(f"No {brand} benchmark available.")
                continue
            cohort_df["Rank"] = cohort_df["ConversionPct_2025"].rank(method="min", ascending=False)
            rank = int(cohort_df.loc[cohort_df["OEM"] == brand, "Rank"].iloc[0])
            leader = cohort_df.sort_values("ConversionPct_2025", ascending=False).iloc[0]
            gap = row["ConversionPct_2025"] - leader["ConversionPct_2025"]
            vts_gap = row["Visits to Sale 2025"] - leader["Visits to Sale 2025"]
            benchmark_card(
                f"{brand} benchmark",
                f"{brand} ranks #{rank} of {len(cohort_df)} in its benchmark set. The Website-to-Contract Conversion Rate gap to {leader['OEM']} is {gap:+.2f}pp. Visits-to-sale gap is {vts_gap:+.0f}.",
                f"Website-to-Contract Conversion Rate: {row['ConversionPct_2025']:.2f}%",
            )


def render_market_weakness(data):
    section("Market Gap Analysis — selected OEM Website-to-Contract performance")
    rows = []
    for market in ["UK", "France", "Germany", "Italy", "Spain"]:
        yoy = yoy_table(data, market, None)
        if yoy.empty:
            continue
        for brand, cohort in [("Toyota", TOYOTA_SET), ("Lexus", LEXUS_SET)]:
            row = get_row(yoy, brand)
            cohort_df = yoy[yoy["OEM"].isin(cohort)].copy()
            if row is None or cohort_df.empty:
                continue
            leader = cohort_df.sort_values("ConversionPct_2025", ascending=False).iloc[0]
            rows.append(
                {
                    "Brand": brand,
                    "Market": market,
                    "Visitor YoY vs Prev Period": row["Visitors YoY %"],
                    "Sales YoY vs Previous Period": row["Sales YoY %"],
                    "W2C Conv Rate": f"{row['ConversionPct_2025']:.2f}%",
                    "Benchmark Leader": leader["OEM"],
                    "Leader W2C Rate": f"{leader['ConversionPct_2025']:.2f}%",
                    "Gap to Leader": row["ConversionPct_2025"] - leader["ConversionPct_2025"],
                }
            )
    if not rows:
        st.info("No selected OEM market gap data available.")
        return

    df = pd.DataFrame(rows).sort_values("Gap to Leader")
    display = df[
        [
            "Brand",
            "Market",
            "Visitor YoY vs Prev Period",
            "Sales YoY vs Previous Period",
            "W2C Conv Rate",
            "Benchmark Leader",
            "Leader W2C Rate",
            "Gap to Leader",
        ]
    ].copy()

    display["Visitor YoY vs Prev Period"] = display["Visitor YoY vs Prev Period"].map(fmt_pct)
    display["Sales YoY vs Previous Period"] = display["Sales YoY vs Previous Period"].map(fmt_pct)
    display["Gap to Leader"] = display["Gap to Leader"].map(fmt_pp)

    badge_cols = ["Visitor YoY vs Prev Period", "Sales YoY vs Previous Period", "Gap to Leader"]
    styler = display.style.map(badge_style, subset=badge_cols)
    st.dataframe(styler, use_container_width=True, hide_index=True)


def make_insight(kind, title, copy, metric, tag, market=None):
    label = {"opportunity": "Opportunity", "risk": "Risk Alert", "intelligence": "Intelligence"}[kind]
    flag = f"<div class='market-flag'>{market}</div>" if market else ""
    return (
        f"<div class='insight-card {kind}'>"
        f"{flag}"
        f"<div class='insight-label {kind}'>{label}</div>"
        f"<div class='insight-title'>{title}</div>"
        f"<div class='insight-copy'>{copy}</div>"
        f"<div class='insight-metric'>{metric}</div>"
        f"<div class='tag'>{tag}</div>"
        f"</div>"
    )


def market_action_text(brand, market, row, leader, gap):
    if market == "Germany":
        return f"{brand} is materially behind {leader['OEM']} in Germany. Prioritise retailer hand-off quality, stock availability messaging and high-intent model pages; this market needs fewer generic visits and more purchase-ready journeys."
    if market == "UK":
        return f"In the UK, {brand} needs sharper lower-funnel progression. Focus on enquiry forms, test-drive booking, finance CTAs and stock routing where intent is most likely leaking before contract."
    if market == "France":
        return f"France needs stronger localisation of offers and dealer follow-up prompts. The opportunity is to turn browsing into configured, dealer-ready demand."
    if market == "Italy":
        return f"Italy looks like a journey-efficiency problem. Reduce friction between model discovery and contact request; simplify price, stock and finance visibility."
    if market == "Spain":
        return f"Spain is relatively more resilient. Protect momentum by scaling the highest-converting traffic sources and making offer, stock and lead response paths more prominent."
    return f"{brand} trails {leader['OEM']} by {gap:+.2f}pp. Prioritise highest-intent contract funnel steps and remove friction from enquiry, finance and stock discovery."


def render_toyota_lexus_recommendations(data):
    section("AI insights — selected OEM market recommendations")
    cards = []
    for market in ["UK", "France", "Germany", "Italy", "Spain"]:
        yoy = yoy_table(data, market, None)
        if yoy.empty:
            continue
        for brand, cohort in [("Toyota", TOYOTA_SET), ("Lexus", LEXUS_SET)]:
            row = get_row(yoy, brand)
            cohort_df = yoy[yoy["OEM"].isin(cohort)].copy()
            if row is None or cohort_df.empty:
                continue
            leader = cohort_df.sort_values("ConversionPct_2025", ascending=False).iloc[0]
            gap = row["ConversionPct_2025"] - leader["ConversionPct_2025"]
            severity = abs(gap) + (0.5 if row["Sales YoY %"] < 0 else 0)
            kind = "risk" if gap < -0.75 or row["Sales YoY %"] < 0 else "opportunity"
            copy = (
                f"{brand} converts at {row['ConversionPct_2025']:.2f}% versus {leader['OEM']} at {leader['ConversionPct_2025']:.2f}%. "
                f"Sales are {fmt_pct(row['Sales YoY %'])} vs previous period; visitors are {fmt_pct(row['Visitors YoY %'])} vs previous period. "
                + market_action_text(brand, market, row, leader, gap)
            )
            cards.append(
                (
                    severity,
                    make_insight(
                        kind,
                        f"{brand}: {market} conversion gap to {leader['OEM']}",
                        copy,
                        f"{gap:+.2f}pp gap",
                        f"{brand} recommendation",
                        market,
                    ),
                )
            )
    cards = [c for _, c in sorted(cards, key=lambda x: x[0], reverse=True)[:8]]
    for i in range(0, len(cards), 3):
        cols = st.columns(3)
        for col, card in zip(cols, cards[i:i + 3]):
            with col:
                st.markdown(card, unsafe_allow_html=True)


def render_top10_table(data, market, selected_oems):
    section("Top 10 brands at a glance")
    yoy = yoy_table(data, market, selected_oems)
    if yoy.empty:
        st.info("No Top 10 data available for this selection.")
        return
    top = yoy.sort_values("UniqueVisitors_2025", ascending=False).head(10).copy()
    rows = []
    for i, (_, r) in enumerate(top.iterrows(), start=1):
        rows.append(
            {
                "#": i,
                "Brand": r["OEM"],
                "Category": cluster_for_oem(r["OEM"]),
                f"Visits current period": fmt_short(r["UniqueVisitors_2025"]),
                f"Visits previous period": fmt_short(r["UniqueVisitors_2024"]),
                "Visitor YoY vs previous period": fmt_pct(r["Visitors YoY %"]),
                "Passenger sales": fmt_int(r["Sales_2025"]),
                "W2C rate": f"{r['ConversionPct_2025']:.2f}%",
                "W2C var vs previous period": fmt_pp(r["Conv Var pp"]),
            }
        )
    display = pd.DataFrame(rows)
    st.dataframe(
        display.style.map(badge_style, subset=["Visitor YoY vs previous period", "W2C var vs previous period"]),
        use_container_width=True,
        hide_index=True,
    )


def build_yoy_growth_chart(data, market, selected_oems, top_n=10):
    yoy = yoy_table(data, market, selected_oems)
    fig = go.Figure()
    if yoy.empty:
        return fig

    metric = "Visitors YoY %"
    yoy = yoy.dropna(subset=[metric])
    top = yoy.sort_values(metric, ascending=False).head(top_n)
    bottom = yoy.sort_values(metric, ascending=True).head(top_n)
    chart_df = pd.concat([top, bottom]).drop_duplicates("OEM").sort_values(metric, ascending=True)
    colors = chart_df[metric].apply(lambda x: "#5ED6AC" if x >= 0 else "#FF5C8A")

    fig.add_trace(
        go.Bar(
            x=chart_df[metric],
            y=chart_df["OEM"],
            orientation="h",
            marker=dict(color=colors),
            text=chart_df[metric].map(lambda x: f"{x:+.1f}%"),
            textposition="outside",
            hovertemplate="<b>%{y}</b><br>Visitor YoY: %{x:+.1f}%<extra></extra>",
        )
    )
    fig.update_layout(
        title=dict(text=f"YoY visit growth — 2025 vs previous period (top & bottom {top_n})", x=0.02, font=dict(size=16, color="#8D96A0")),
        height=620,
        margin=dict(l=120, r=80, t=70, b=50),
        plot_bgcolor=WHITE,
        paper_bgcolor=WHITE,
        showlegend=False,
    )
    min_x = min(chart_df[metric].min() * 1.25, -5)
    max_x = max(chart_df[metric].max() * 1.25, 5)
    fig.update_xaxes(range=[min_x, max_x], ticksuffix="%", gridcolor="#E2E6EA", zeroline=True, zerolinecolor="#B8C0C8")
    fig.update_yaxes(gridcolor="#FFFFFF")
    return fig


def render_market_insights(data, market, selected_oems):
    section("AI insights — market performance")
    yoy = yoy_table(data, market, selected_oems)
    if yoy.empty:
        st.info("No insight data available.")
        return

    top_conv = yoy.sort_values("ConversionPct_2025", ascending=False).iloc[0]
    top_sales = yoy.sort_values("Sales YoY %", ascending=False).iloc[0]
    top_visits = yoy.sort_values("Visitors YoY %", ascending=False).iloc[0]
    worst_conv = yoy.sort_values("Conv Var pp", ascending=True).iloc[0]
    weakest = yoy.sort_values("ConversionPct_2025", ascending=True).iloc[0]

    cards = [
        make_insight("opportunity", f"{top_conv['OEM']}: strongest Website-to-Contract performer", f"{top_conv['OEM']} leads this cohort on Website-to-Contract Conversion Rate. Visits-to-sale is {top_conv['Visits to Sale 2025']:.0f}, showing stronger website-to-contract efficiency than peers.", f"{top_conv['ConversionPct_2025']:.2f}% conv", top_conv["OEM"]),
        make_insight("opportunity", f"{top_sales['OEM']}: strongest sales growth", f"{top_sales['OEM']} posted the strongest passenger sales growth, with sales {fmt_pct(top_sales['Sales YoY %'])} vs previous period and visitors {fmt_pct(top_sales['Visitors YoY %'])} vs previous period.", f"{fmt_pct(top_sales['Sales YoY %'])} sales", top_sales["OEM"]),
        make_insight("intelligence", f"{top_visits['OEM']}: fastest visitor growth", f"{top_visits['OEM']} generated the strongest visitor growth. The key question is whether this awareness converts into contracts at the same rate.", f"{fmt_pct(top_visits['Visitors YoY %'])} visits", top_visits["OEM"]),
        make_insight("risk", f"{worst_conv['OEM']}: Website-to-Contract movement deteriorating", f"{worst_conv['OEM']} had the weakest conversion movement: {fmt_pp(worst_conv['Conv Var pp'])} vs previous period. This suggests lower-quality traffic or leakage deeper in the funnel.", f"{fmt_pp(worst_conv['Conv Var pp'])}", worst_conv["OEM"]),
        make_insight("risk", f"{weakest['OEM']}: weakest Website-to-Contract Conversion Rate", f"{weakest['OEM']} has the lowest Website-to-Contract Conversion Rate rate in this selection, creating a clear efficiency gap against the cohort leader.", f"{weakest['ConversionPct_2025']:.2f}% conv", weakest["OEM"]),
    ]
    for i in range(0, len(cards), 3):
        cols = st.columns(3)
        for col, card in zip(cols, cards[i:i + 3]):
            with col:
                st.markdown(card, unsafe_allow_html=True)


def get_brandfetch_client_id():
    try:
        return st.secrets.get("BRANDFETCH_CLIENT_ID", "")
    except Exception:
        return ""


def get_logo_url(oem):
    client_id = get_brandfetch_client_id()
    domain = LOGO_DOMAIN_MAP.get(oem)
    if not client_id or not domain:
        return None
    return f"https://cdn.brandfetch.io/{domain}?c={client_id}"


def add_logo_images(fig, chart_df, x_max, y_max):
    for _, row in chart_df.iterrows():
        logo = get_logo_url(row["OEM"])
        if not logo:
            continue
        fig.add_layout_image(
            dict(
                source=logo,
                xref="x",
                yref="y",
                x=row["UniqueVisitors"],
                y=row["ConversionPct"],
                sizex=x_max * 0.035,
                sizey=y_max * 0.055,
                xanchor="center",
                yanchor="middle",
                sizing="contain",
                opacity=0.95,
                layer="above",
            )
        )


def build_bubble_chart(chart_df, selected_oems, market, year_view, show_logos, x_axis_metric="Unique Visitors", view_label="OEM"):
    fig = go.Figure()
    if chart_df.empty:
        return fig
    x_col = "Sales" if x_axis_metric == "Passenger Car Sales" else "UniqueVisitors"
    size_col = "UniqueVisitors" if x_col == "Sales" else "Sales"
    x_title = "Passenger Car Sales" if x_col == "Sales" else "Unique visitors"
    size_title = "Unique visitors" if size_col == "UniqueVisitors" else "Passenger Car Sales"
    max_size = max(chart_df[size_col].max(), 1)
    max_x = max(chart_df[x_col].max() * 1.18, 1)
    max_y = max(chart_df["ConversionPct"].max() * 1.25, 0.1)
    current = chart_df[chart_df["Year"] == 2025].copy()
    x_mid = current[x_col].median() if not current.empty else chart_df[x_col].median()
    y_mid = current["ConversionPct"].median() if not current.empty else chart_df["ConversionPct"].median()
    for x0,x1,y0,y1,fill in [(0,x_mid,y_mid,max_y,"rgba(0,159,227,0.05)"),(x_mid,max_x,y_mid,max_y,"rgba(18,199,107,0.07)"),(0,x_mid,0,y_mid,"rgba(111,111,111,0.04)"),(x_mid,max_x,0,y_mid,"rgba(255,176,0,0.06)")]:
        fig.add_shape(type="rect", x0=x0, x1=x1, y0=y0, y1=y1, fillcolor=fill, line=dict(width=0), layer="below")
    fig.add_shape(type="line", x0=x_mid, x1=x_mid, y0=0, y1=max_y, line=dict(color="rgba(10,35,66,0.12)", width=1), layer="below")
    fig.add_shape(type="line", x0=0, x1=max_x, y0=y_mid, y1=y_mid, line=dict(color="rgba(10,35,66,0.12)", width=1), layer="below")
    id_col = "DisplayName" if "DisplayName" in chart_df.columns else "OEM"
    if year_view == "Previous and current + shift":
        for item in chart_df[id_col].dropna().unique():
            d24 = chart_df[(chart_df[id_col] == item) & (chart_df["Year"] == 2024)]
            d25 = chart_df[(chart_df[id_col] == item) & (chart_df["Year"] == 2025)]
            if len(d24) and len(d25):
                r24, r25 = d24.iloc[0], d25.iloc[0]
                fig.add_annotation(x=r25[x_col], y=r25["ConversionPct"], ax=r24[x_col], ay=r24["ConversionPct"], xref="x", yref="y", axref="x", ayref="y", showarrow=True, arrowhead=3, arrowsize=1.1, arrowwidth=1.4, arrowcolor="#6F7782", opacity=0.70)
        df_prev = chart_df[chart_df["Year"] == 2024].copy()
        if not df_prev.empty:
            fig.add_trace(go.Scatter(x=df_prev[x_col], y=df_prev["ConversionPct"], mode="markers", name=PREVIOUS_LABEL, marker=dict(size=df_prev[size_col].apply(lambda s: max(16, math.sqrt(s / max_size) * 68)), color="rgba(190,198,208,0.45)", line=dict(width=1.2, color="rgba(10,35,66,0.25)")), customdata=df_prev[[id_col,"Market","Sales","UniqueVisitors","ConversionPct"]], hovertemplate="<b>%{customdata[0]}</b><br>Market: %{customdata[1]}<br>Sales: %{customdata[2]:,.0f}<br>Visitors: %{customdata[3]:,.0f}<br>W2C: %{customdata[4]:.2f}%<extra></extra>"))
    primary_year = 2025 if year_view == "Previous and current + shift" else (2024 if year_view == "Previous period" else 2025)
    primary = chart_df[chart_df["Year"] == primary_year].copy()
    if not primary.empty:
        if "Cluster" not in primary.columns:
            primary["Cluster"] = primary["OEM"].map(cluster_for_oem) if "OEM" in primary.columns else "MM5 Markets"
        for cluster in sorted(primary["Cluster"].astype(str).unique()):
            dfc = primary[primary["Cluster"].astype(str) == cluster]
            color = CLUSTER_COLORS.get(cluster, "#009FE3")
            fig.add_trace(go.Scatter(x=dfc[x_col], y=dfc["ConversionPct"], mode="markers+text" if not show_logos else "markers", text=None if show_logos else dfc[id_col], textposition="top center", name=f"{cluster} · {CURRENT_LABEL}", marker=dict(size=dfc[size_col].apply(lambda s: max(24, math.sqrt(s / max_size) * 86)), color=color, opacity=0.78, line=dict(width=2.2, color="white")), customdata=dfc[[id_col,"Market","Sales","UniqueVisitors","ConversionPct"]], hovertemplate="<b>%{customdata[0]}</b><br>Market: %{customdata[1]}<br>Sales: %{customdata[2]:,.0f}<br>Visitors: %{customdata[3]:,.0f}<br>W2C: %{customdata[4]:.2f}%<extra></extra>"))
    if show_logos and view_label == "OEM" and not primary.empty:
        add_logo_images(fig, primary.rename(columns={id_col:"OEM"}), max_x, max_y)
    fig.update_layout(title=dict(text=f"{market} | {x_title} vs Website-to-Contract Conversion Rate", x=0.02, y=0.98, font=dict(size=18, color="#0A2342")), xaxis_title=x_title, yaxis_title="Website-to-Contract Conversion Rate", height=650, hovermode="closest", legend=dict(orientation="h", y=1.08, x=0, font=dict(size=11)), margin=dict(l=70, r=35, t=95, b=65), plot_bgcolor="#FFFFFF", paper_bgcolor="#FFFFFF")
    fig.update_xaxes(range=[0, max_x], tickformat=",", gridcolor="#EDF1F5", zeroline=False, showline=True, linecolor="#E1E7EF")
    fig.update_yaxes(range=[0, max_y], ticksuffix="%", gridcolor="#EDF1F5", zeroline=False, showline=True, linecolor="#E1E7EF")
    return fig


def assistant_card(title, copy, metric, tag, style="primary"):
    return (
        f"<div class='assistant-card {style}'>"
        f"<div class='assistant-card-title'>{title}</div>"
        f"<div class='assistant-card-copy'>{copy}</div>"
        f"<div class='assistant-card-metric'>{metric}</div>"
        f"<div class='assistant-card-tag'>{tag}</div>"
        f"</div>"
    )


def extract_market_from_question(question):
    q = question.lower()
    market_aliases = {
        "mm5": "MM5",
        "uk": "UK",
        "united kingdom": "UK",
        "france": "France",
        "french": "France",
        "germany": "Germany",
        "german": "Germany",
        "italy": "Italy",
        "italian": "Italy",
        "spain": "Spain",
        "spanish": "Spain",
    }
    for key, market in market_aliases.items():
        if key in q:
            return market
    return "MM5"


def extract_brand_from_question(question, available_oems):
    q = question.lower()

    aliases = {
        "volkswagen": "VW",
        "vw": "VW",
        "byd": "BYD Auto",
        "byd auto": "BYD Auto",
        "ds automobiles": "DS",
        "citroen": "Citroen",
    }

    # Match aliases and brand names as whole tokens/phrases only.
    # This prevents short brands such as "DS" being detected inside words like "leads".
    candidates = []
    for brand in available_oems:
        candidates.append((brand.lower(), brand))
    for alias, brand in aliases.items():
        if brand in available_oems:
            candidates.append((alias.lower(), brand))

    for token, brand in sorted(candidates, key=lambda x: len(x[0]), reverse=True):
        pattern = r"(?<![a-z0-9])" + re.escape(token) + r"(?![a-z0-9])"
        if re.search(pattern, q):
            return brand

    return None


def assistant_result_html(cards):
    return "<div class='assistant-result-grid'>" + "".join(cards) + "</div>"


def top_brands_cards(data, market, selected_oems, question):
    q = question.lower()
    yoy = yoy_table(data, market, selected_oems)
    if yoy.empty:
        return [assistant_card("No data available", f"No comparable data was found for {market}.", "—", market, "neutral")]

    if "sales" in q:
        metric = "Sales_2025"
        label = "passenger sales"
        formatter = fmt_metric_number
    elif "visitor" in q or "traffic" in q or "audience" in q:
        metric = "UniqueVisitors_2025"
        label = "unique visitors"
        formatter = fmt_metric_number
    else:
        metric = "ConversionPct_2025"
        label = "Website-to-Contract Conversion Rate"
        formatter = lambda x: f"{x:.2f}%"

    ascending = any(term in q for term in ["worst", "weakest", "lowest", "underperform"])
    ranked = yoy.sort_values(metric, ascending=ascending).head(5)

    cards = []
    for idx, (_, row) in enumerate(ranked.iterrows(), start=1):
        style = "risk" if ascending else ("positive" if idx == 1 else "primary")
        direction = "lowest" if ascending else "highest"
        cards.append(
            assistant_card(
                f"#{idx}: {row['OEM']}",
                f"In {market}, {row['OEM']} is ranked #{idx} for {label} in the selected set. "
                f"Sales changed {fmt_pct(row['Sales YoY %'])}; unique visitors changed {fmt_pct(row['Visitors YoY %'])} versus {PREVIOUS_LABEL}.",
                formatter(row[metric]),
                f"{direction} {label}",
                style,
            )
        )

    leader = ranked.iloc[0]
    cards.append(
        assistant_card(
            "Read this correctly",
            f"This is not proof that {leader['OEM']} has the best website. W2C performance can reflect brand demand, product fit, pricing, stock, retailer follow-up and network effects.",
            "macro funnel",
            f"{CURRENT_LABEL} vs {PREVIOUS_LABEL}",
            "neutral",
        )
    )

    return cards


def brand_cards(data, market, brand):
    yoy = yoy_table(data, market, None)
    if yoy.empty:
        return [assistant_card("No data available", f"No comparable data was found for {market}.", "—", market, "neutral")]

    row = get_row(yoy, brand)
    if row is None:
        return [assistant_card("Brand not found", f"I could not find {brand} in {market} for the active comparison.", "—", brand, "neutral")]

    cohort = sorted(yoy["OEM"].unique().tolist())
    cohort_df = yoy[yoy["OEM"].isin(cohort)].copy()
    leader = cohort_df.sort_values("ConversionPct_2025", ascending=False).iloc[0] if not cohort_df.empty else yoy.sort_values("ConversionPct_2025", ascending=False).iloc[0]
    gap = row["ConversionPct_2025"] - leader["ConversionPct_2025"]
    rank = int(cohort_df["ConversionPct_2025"].rank(method="min", ascending=False).loc[row.name]) if row.name in cohort_df.index else int(yoy["ConversionPct_2025"].rank(method="min", ascending=False).loc[row.name])

    visitor_style = "positive" if row["Visitors YoY %"] >= 0 else "risk"
    sales_style = "positive" if row["Sales YoY %"] >= 0 else "risk"
    conv_style = "positive" if row["Conv Var pp"] >= 0 else "risk"

    cards = [
        assistant_card(
            f"{brand} W2C performance in {market}",
            f"{brand}'s {CURRENT_LABEL} Website-to-Contract Conversion Rate is ranked #{rank} within the relevant comparison set.",
            f"{row['ConversionPct_2025']:.2f}%",
            "W2C rate",
            "primary",
        ),
        assistant_card(
            "Unique visitor movement",
            f"Unique visitor demand moved versus {PREVIOUS_LABEL}. This helps separate demand generation from conversion efficiency.",
            fmt_pct(row["Visitors YoY %"]),
            fmt_metric_number(row["UniqueVisitors_2025"]),
            visitor_style,
        ),
        assistant_card(
            "Passenger sales movement",
            f"Passenger sales moved versus {PREVIOUS_LABEL}. Read this alongside visitor movement to understand whether demand is translating into contracts.",
            fmt_pct(row["Sales YoY %"]),
            fmt_metric_number(row["Sales_2025"]),
            sales_style,
        ),
        assistant_card(
            "W2C rate movement",
            f"The W2C rate movement indicates whether conversion efficiency improved or deteriorated versus the prior period.",
            fmt_pp(row["Conv Var pp"]),
            f"vs {PREVIOUS_LABEL}",
            conv_style,
        ),
        assistant_card(
            "Benchmark leader gap",
            f"The benchmark leader is {leader['OEM']} at {leader['ConversionPct_2025']:.2f}%.",
            fmt_pp(gap),
            f"gap to {leader['OEM']}",
            "positive" if gap >= 0 else "risk",
        ),
        assistant_card(
            "Likely diagnostic route",
            "Use this as a prompt for deeper diagnosis across marketing quality, model-page content, offers, stock visibility, lead handling and retailer follow-up.",
            "next step",
            "diagnostic",
            "neutral",
        ),
    ]
    return cards


def toyota_lexus_cards(data, market):
    yoy = yoy_table(data, market, None)
    if yoy.empty:
        return [assistant_card("No data available", f"No comparable data was found for {market}.", "—", market, "neutral")]

    cards = []
    for brand in ["Toyota", "Lexus"]:
        row = get_row(yoy, brand)
        if row is None:
            continue
        cohort = TOYOTA_SET if brand == "Toyota" else LEXUS_SET
        cohort_df = yoy[yoy["OEM"].isin(cohort)].copy()
        leader = cohort_df.sort_values("ConversionPct_2025", ascending=False).iloc[0]
        gap = row["ConversionPct_2025"] - leader["ConversionPct_2025"]
        cards.extend(
            [
                assistant_card(
                    f"{brand}: W2C rate in {market}",
                    f"{brand}'s {CURRENT_LABEL} Website-to-Contract Conversion Rate is compared against its benchmark cohort.",
                    f"{row['ConversionPct_2025']:.2f}%",
                    brand,
                    "primary",
                ),
                assistant_card(
                    f"{brand}: benchmark gap",
                    f"The benchmark leader is {leader['OEM']} at {leader['ConversionPct_2025']:.2f}%.",
                    fmt_pp(gap),
                    f"leader: {leader['OEM']}",
                    "positive" if gap >= 0 else "risk",
                ),
                assistant_card(
                    f"{brand}: demand vs sales movement",
                    f"Unique visitors moved {fmt_pct(row['Visitors YoY %'])}; sales moved {fmt_pct(row['Sales YoY %'])} versus {PREVIOUS_LABEL}.",
                    f"{fmt_pct(row['Visitors YoY %'])} visitors",
                    f"{fmt_pct(row['Sales YoY %'])} sales",
                    "positive" if row["Conv Var pp"] >= 0 else "risk",
                ),
            ]
        )

    cards.append(
        assistant_card(
            "How to use this",
            "Where visitor growth is positive but sales or W2C decline, the likely issue is not just digital reach. Investigate traffic quality, offers, stock, retailer execution and proposition fit.",
            "diagnostic lens",
            "Selected OEM",
            "neutral",
        )
    )

    return cards


def explain_metric_cards():
    return [
        assistant_card(
            "Website Conversion Rate",
            "This is the existing website effectiveness metric used by the group today. It measures visitors who complete a lead or online action divided by total website visitors.",
            "converted visitors / total visitors",
            "website effectiveness",
            "neutral",
        ),
        assistant_card(
            "Website-to-Contract Conversion Rate",
            "This dashboard uses a broader commercial efficiency metric: passenger new car Passenger Car Sales divided by unique website visitors.",
            "contracts / unique visitors",
            "macro conversion",
            "primary",
        ),
        assistant_card(
            "Why the distinction matters",
            "A low W2C rate does not automatically mean a weak website. It can reflect brand demand, pricing, product, stock, retailer follow-up, channel mix or market structure.",
            "not website-only",
            "interpretation",
            "positive",
        ),
        assistant_card(
            "Unique visitor source",
            "Unique visitors are sourced from Similarweb and should be used as a consistent market benchmark, not as a replacement for first-party web analytics.",
            "Similarweb",
            "data source",
            "neutral",
        ),
        assistant_card(
            "Sales source",
            "Passenger car sales are sourced from Marklines. The metric does not include every possible sales channel interpretation such as fleet or tactical registrations.",
            "Marklines",
            "sales source",
            "neutral",
        ),
        assistant_card(
            "Best use",
            "Use W2C to provoke bigger-picture questions about why demand does or does not translate into Passenger Car Sales.",
            "diagnostic",
            "use case",
            "primary",
        ),
    ]


def generate_assistant_cards(data, question, selected_oems):
    q = question.lower().strip()
    market = extract_market_from_question(q)
    brand = extract_brand_from_question(q, sorted(data["OEM"].unique()))

    if not q:
        return []

    if any(term in q for term in ["define", "definition", "methodology", "what is", "explain"]):
        if "conversion" in q or "w2c" in q or "website" in q:
            return explain_metric_cards()

    # Leaderboard-style questions must be answered before brand detection.
    # Example: "who leads W2C in Germany?" should return the #1 OEM, not detect "DS" inside "leads".
    leaderboard_terms = ["leader", "leads", "lead", "best", "top", "outperform", "highest", "worst", "weakest", "lowest", "underperform"]
    metric_terms = ["w2c", "conversion", "sales", "visitor", "traffic", "audience", "contract"]
    if any(term in q for term in leaderboard_terms) or any(term in q for term in metric_terms):
        if brand and any(term in q for term in ["how is", "how are", "performance", "performing", "gap", "compare"]):
            return brand_cards(data, market, brand)

        if any(term in q for term in leaderboard_terms) or any(term in q for term in ["highest", "lowest", "top", "best", "worst"]):
            return top_brands_cards(data, market, selected_oems, q)

    if brand:
        return brand_cards(data, market, brand)

    if any(term in q for term in ["sales", "visitor", "traffic", "audience"]):
        return top_brands_cards(data, market, selected_oems, q)

    return [
        assistant_card(
            "Try asking a more specific data question",
            "I can answer questions about leaders, weakest performers, competitor gaps, brand performance, sales, visitors and W2C conversion by market.",
            "Example",
            "Who leads W2C in Germany?",
            "neutral",
        ),
        assistant_card(
            "Useful prompts",
            "Try: 'How is Honda performing in Germany?', 'Who has the highest W2C rate in MM5?', or 'Explain the conversion metric'.",
            "3 prompts",
            "guided assistant",
            "primary",
        ),
    ]


def render_data_assistant(data, selected_oems):
    st.markdown(
        "<div class='da-simple'>"
        "<div class='da-simple-title'>Welcome. How can I help?</div>"
        "<div class='da-simple-copy'>Ask a question about markets, OEMs, Website-to-Contract Conversion Rate, sales, visitors or competitor gaps. Responses are generated from the dashboard dataset only.</div>"
        "<div class='da-simple-label'>Example questions</div>"
        "<div class='da-simple-questions'>"
        "<div class='da-simple-question'>Who leads W2C in Germany?</div>"
        "<div class='da-simple-question'>How is Honda performing in France?</div>"
        "<div class='da-simple-question'>Compare Honda and Volkswagen in MM5</div>"
                "</div>"
        "</div>",
        unsafe_allow_html=True,
    )

    st.markdown(
        f"<div class='da-context-bar'>"
        f"<div></div>"
        f"<div class='da-context-chip'>{CURRENT_LABEL} vs {PREVIOUS_LABEL}</div>"
        f"</div>",
        unsafe_allow_html=True,
    )

    question = st.text_input(
        "Ask the data",
        placeholder="Search...",
        key="data_assistant_question",
    )

    if question:
        cards = generate_assistant_cards(data, question, selected_oems)[:3]
        st.markdown(assistant_result_html(cards), unsafe_allow_html=True)
    else:
        st.info("Start by typing a question above. Example: 'Which OEM has the highest W2C rate in the UK?'")



def ppt_safe_text(value):
    if value is None:
        return ""
    return str(value).replace("&", "and")


def selected_oems_from_report_filters(data, categories, preset, picked_oems):
    if preset and preset != "None":
        oems = resolve_tme_preset(data, preset)
    elif categories:
        oems = oems_for_categories(data, categories)
    else:
        oems = sorted(data["OEM"].unique())
    if picked_oems:
        picked = set(picked_oems)
        oems = [o for o in oems if o in picked]
    return sorted(oems)


def ppt_add_title(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.45), Inches(0.30), Inches(12.4), Inches(0.75))
    p = tx.text_frame.paragraphs[0]
    p.text = ppt_safe_text(title)
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 35, 66)
    if subtitle:
        stx = slide.shapes.add_textbox(Inches(0.48), Inches(0.95), Inches(12.0), Inches(0.35))
        sp = stx.text_frame.paragraphs[0]
        sp.text = ppt_safe_text(subtitle)
        sp.font.size = Pt(11)
        sp.font.color.rgb = RGBColor(111, 119, 130)


def ppt_add_kpi(slide, x, y, label, value, delta=None):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.75), Inches(1.25))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(247, 249, 252)
    box.line.color.rgb = RGBColor(230, 233, 237)
    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = ppt_safe_text(label)
    p1.font.size = Pt(9)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(111, 119, 130)
    p2 = tf.add_paragraph()
    p2.text = ppt_safe_text(value)
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0, 0, 0)
    if delta:
        p3 = tf.add_paragraph()
        p3.text = ppt_safe_text(delta)
        p3.font.size = Pt(9)
        p3.font.bold = True
        p3.font.color.rgb = RGBColor(18, 199, 107) if str(delta).startswith("+") else RGBColor(255, 47, 109) if str(delta).startswith("-") else RGBColor(111, 119, 130)


def ppt_add_table(slide, df, x, y, w, h, font_size=8):
    if df.empty:
        return
    rows, cols = min(len(df) + 1, 14), len(df.columns)
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for c, col in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = ppt_safe_text(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(247, 249, 252)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = RGBColor(10, 35, 66)
    for r_idx, (_, row) in enumerate(df.head(rows - 1).iterrows(), start=1):
        for c, col in enumerate(df.columns):
            cell = table.cell(r_idx, c)
            cell.text = ppt_safe_text(row[col])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(48, 54, 66)


def ppt_insight_bullets(data, market, oems):
    yoy = yoy_table(data, market, oems)
    if yoy.empty:
        return ["No comparable data available for the selected source."]
    top_w2c = yoy.sort_values("ConversionPct_2025", ascending=False).iloc[0]
    top_sales = yoy.sort_values("Sales_2025", ascending=False).iloc[0]
    fastest_visitors = yoy.sort_values("Visitors YoY %", ascending=False).iloc[0]
    weakest_conv = yoy.sort_values("Conv Var pp", ascending=True).iloc[0]
    bullets = [
        f"{top_w2c['OEM']} leads W2C efficiency at {top_w2c['ConversionPct_2025']:.2f}%.",
        f"{top_sales['OEM']} has the largest Passenger Car Sales volume at {fmt_metric_number(top_sales['Sales_2025'])}.",
        f"{fastest_visitors['OEM']} has the strongest visitor momentum at {fmt_pct(fastest_visitors['Visitors YoY %'])}.",
        f"{weakest_conv['OEM']} shows the weakest W2C movement at {fmt_pp(weakest_conv['Conv Var pp'])}.",
    ]
    tl = yoy[yoy["OEM"].isin(["Toyota", "Lexus"])]
    if not tl.empty:
        tl_names = ", ".join(tl["OEM"].tolist())
        tl_avg = tl["ConversionPct_2025"].mean()
        bullets.append(f"{tl_names} average W2C rate is {tl_avg:.2f}% in the selected comparison set.")
    return bullets[:5]


def create_insight_report_ppt(data, market, categories, preset, picked_oems, report_style):
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Add python-pptx to requirements.txt and redeploy.")

    oems = selected_oems_from_report_filters(data, categories, preset, picked_oems)
    if not oems:
        oems = sorted(data["OEM"].unique())

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    yoy = yoy_table(data, market, oems)
    summary = market_kpi_summary(data, market, oems)

    # Slide 1
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
    bg.line.color.rgb = RGBColor(0, 0, 0)
    title = slide.shapes.add_textbox(Inches(0.65), Inches(0.8), Inches(11.6), Inches(1.6))
    p = title.text_frame.paragraphs[0]
    p.text = "OEM Market Intelligence Report"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle = slide.shapes.add_textbox(Inches(0.68), Inches(2.25), Inches(11.3), Inches(0.8))
    sp = subtitle.text_frame.paragraphs[0]
    sp.text = f"{market} | {CURRENT_LABEL} vs {PREVIOUS_LABEL} | {report_style}"
    sp.font.size = Pt(18)
    sp.font.color.rgb = RGBColor(220, 230, 238)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(7.25), prs.slide_width, Inches(0.25))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0, 159, 227)
    bar.line.color.rgb = RGBColor(0, 159, 227)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    ppt_add_title(slide, "Executive snapshot", f"Source: {market}; OEMs selected: {len(oems)}")
    if summary:
        ppt_add_kpi(slide, 0.6, 1.45, "Passenger Car Sales", fmt_metric_number(summary["sales"]), fmt_pct(summary["sales_delta"]) if summary["sales_delta"] is not None else None)
        ppt_add_kpi(slide, 3.55, 1.45, "Unique visitors", fmt_metric_number(summary["visitors"]), fmt_pct(summary["visitor_delta"]) if summary["visitor_delta"] is not None else None)
        ppt_add_kpi(slide, 6.5, 1.45, "W2C rate", f"{summary['conv']:.2f}%", fmt_pp(summary["conv_delta"]))
        ppt_add_kpi(slide, 9.45, 1.45, "Visits per sale", fmt_int(summary["visits_sale"]) if summary["visits_sale"] else "n/a", "lower is better")
    bullets = ppt_insight_bullets(data, market, oems)
    box = slide.shapes.add_textbox(Inches(0.75), Inches(3.25), Inches(12.0), Inches(3.2))
    tf = box.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {ppt_safe_text(b)}"
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(10, 35, 66)
        p.space_after = Pt(8)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    ppt_add_title(slide, "Leading competitors", "Top performers by Website-to-Contract Conversion Rate")
    if not yoy.empty:
        top = yoy.sort_values("ConversionPct_2025", ascending=False).head(10).copy()
        table_df = pd.DataFrame({
            "Rank": range(1, len(top) + 1),
            "Brand": top["OEM"],
            "Category": top["OEM"].map(cluster_for_oem),
            "W2C": top["ConversionPct_2025"].map(lambda x: f"{x:.2f}%"),
            "Sales": top["Sales_2025"].map(fmt_metric_number),
            "Visitors": top["UniqueVisitors_2025"].map(fmt_metric_number),
            "Visitor YoY": top["Visitors YoY %"].map(fmt_pct),
            "Sales YoY": top["Sales YoY %"].map(fmt_pct),
        })
        ppt_add_table(slide, table_df, 0.45, 1.25, 12.4, 5.55, font_size=7)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    ppt_add_title(slide, "Selected OEM position", "How the selected OEM compares against the selected source")
    if not yoy.empty:
        tl = yoy[yoy["OEM"].isin(["Toyota", "Lexus"])].copy()
        if not tl.empty:
            tl_df = pd.DataFrame({
                "Brand": tl["OEM"],
                "Category": tl["OEM"].map(cluster_for_oem),
                "W2C": tl["ConversionPct_2025"].map(lambda x: f"{x:.2f}%"),
                "W2C Var": tl["Conv Var pp"].map(fmt_pp),
                "Sales": tl["Sales_2025"].map(fmt_metric_number),
                "Sales YoY": tl["Sales YoY %"].map(fmt_pct),
                "Visitors": tl["UniqueVisitors_2025"].map(fmt_metric_number),
                "Visitor YoY": tl["Visitors YoY %"].map(fmt_pct),
            })
            ppt_add_table(slide, tl_df, 0.6, 1.4, 12.1, 1.5, font_size=9)
        insight = "Use this slide to challenge whether selected OEM performance is a demand issue, conversion-quality issue, or wider commercial issue."
        box = slide.shapes.add_textbox(Inches(0.75), Inches(3.45), Inches(11.8), Inches(1.6))
        p = box.text_frame.paragraphs[0]
        p.text = insight
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(10, 35, 66)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    ppt_add_title(slide, "Category view", "Which OEM groups are driving the selected market source")
    if not yoy.empty:
        cat = yoy.copy()
        cat["Category"] = cat["OEM"].map(cluster_for_oem)
        grouped = cat.groupby("Category", dropna=False).agg(
            Sales=("Sales_2025", "sum"),
            Visitors=("UniqueVisitors_2025", "sum"),
        ).reset_index()
        grouped["W2C"] = grouped["Sales"] / grouped["Visitors"] * 100
        grouped = grouped.sort_values("W2C", ascending=False)
        table_df = pd.DataFrame({
            "Category": grouped["Category"],
            "W2C": grouped["W2C"].map(lambda x: f"{x:.2f}%"),
            "Sales": grouped["Sales"].map(fmt_metric_number),
            "Visitors": grouped["Visitors"].map(fmt_metric_number),
        })
        ppt_add_table(slide, table_df, 0.8, 1.35, 11.7, 4.0, font_size=10)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    ppt_add_title(slide, "Planning questions", "Use these prompts in the QBR discussion")
    questions = [
        "Which front-runner OEMs are gaining both demand and sales efficiency?",
        "Where is visitor growth failing to translate into Passenger Car Sales?",
        "Which competitor categories are creating the strongest pressure?",
        "Where should selected OEM investigate product, offer, stock, retailer or network constraints?",
        "Which signals should shape the next quarterly planning cycle?",
    ]
    qbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.8), Inches(4.8))
    tf = qbox.text_frame
    tf.clear()
    for i, q in enumerate(questions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i+1}. {q}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(10, 35, 66)
        p.space_after = Pt(12)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    ppt_add_title(slide, "Methodology and caveats", "How to read this report")
    notes = [
        "Website-to-Contract Conversion Rate = Passenger Car Sales / unique website visitors.",
        "Unique visitors source: Similarweb. Passenger Car Sales source: Marklines.",
        "This is not a total market-share report. It excludes fleet, LCV and tactical registrations.",
        "Use outputs as market-intelligence prompts, not as proof of direct website causality.",
    ]
    nbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.8), Inches(4.8))
    tf = nbox.text_frame
    tf.clear()
    for i, n in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {ppt_safe_text(n)}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(10, 35, 66)
        p.space_after = Pt(10)

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()


def selected_oem_market_insight_rows(data, brand):
    rows = []
    for market in ["UK", "France", "Germany", "Italy", "Spain"]:
        yoy = yoy_table(data, market, None)
        row = get_row(yoy, brand)
        if row is None or yoy.empty:
            continue
        rank_df = yoy.sort_values("ConversionPct_2025", ascending=False).reset_index(drop=True)
        rank = int(rank_df.index[rank_df["OEM"] == brand][0]) + 1 if brand in set(rank_df["OEM"]) else None
        market_avg = yoy["ConversionPct_2025"].mean()
        leader = rank_df.iloc[0]
        rows.append({
            "Market": market,
            "OEM": brand,
            "W2C": row["ConversionPct_2025"],
            "W2C Var": row["Conv Var pp"],
            "Sales": row["Sales_2025"],
            "Sales YoY": row["Sales YoY %"],
            "Visitors": row["UniqueVisitors_2025"],
            "Visitors YoY": row["Visitors YoY %"],
            "Visits to Sale": row["Visits to Sale 2025"],
            "Rank": rank,
            "Peer Count": len(rank_df),
            "Market Avg W2C": market_avg,
            "Gap to Market Avg": row["ConversionPct_2025"] - market_avg,
            "Leader": leader["OEM"],
            "Leader W2C": leader["ConversionPct_2025"],
            "Gap to Leader": row["ConversionPct_2025"] - leader["ConversionPct_2025"],
        })
    return pd.DataFrame(rows)




def ppt_valtech_header(slide, section_label="OEM INTELLIGENCE"):
    """Add simple Valtech-style branded header/footer elements to a slide."""
    # top wordmark
    tx = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(3.2), Inches(0.35))
    p = tx.text_frame.paragraphs[0]
    p.text = "Valtech ✳"
    p.font.size = Pt(20)
    p.font.bold = False
    p.font.color.rgb = RGBColor(0, 0, 0)

    # red active mark
    mark = slide.shapes.add_shape(1, Inches(0.45), Inches(0.78), Inches(0.28), Inches(0.055))
    mark.fill.solid()
    mark.fill.fore_color.rgb = RGBColor(255, 75, 85)
    mark.line.color.rgb = RGBColor(255, 75, 85)

    lab = slide.shapes.add_textbox(Inches(0.45), Inches(0.95), Inches(4.0), Inches(0.25))
    lp = lab.text_frame.paragraphs[0]
    lp.text = section_label
    lp.font.size = Pt(8)
    lp.font.bold = True
    lp.font.color.rgb = RGBColor(111, 111, 111)

    # bottom blue rule
    bar = slide.shapes.add_shape(1, Inches(0), Inches(7.36), Inches(13.333), Inches(0.14))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0, 60, 179)
    bar.line.color.rgb = RGBColor(0, 60, 179)


def ppt_add_textbox(slide, x, y, w, h, text, size=16, bold=False, color=(0, 0, 0)):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = ppt_safe_text(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    return box


def ppt_add_valtech_card(slide, x, y, w, h, title, value, body=None, accent="blue"):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = RGBColor(0, 0, 0)
    # accent line
    accent_rgb = RGBColor(0, 60, 179) if accent == "blue" else RGBColor(255, 75, 85)
    line = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.07), Inches(h))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_rgb
    line.line.color.rgb = accent_rgb
    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = ppt_safe_text(title)
    p1.font.size = Pt(9)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(111, 111, 111)
    p2 = tf.add_paragraph()
    p2.text = ppt_safe_text(value)
    p2.font.size = Pt(22)
    p2.font.bold = False
    p2.font.color.rgb = RGBColor(0, 0, 0)
    if body:
        p3 = tf.add_paragraph()
        p3.text = ppt_safe_text(body)
        p3.font.size = Pt(9)
        p3.font.color.rgb = RGBColor(111, 111, 111)


def create_oem_insight_storyteller_ppt(data, selected_brand, rows):
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Add python-pptx to requirements.txt and redeploy.")
    if rows.empty:
        raise RuntimeError("No insight rows available to export.")

    best = rows.sort_values("W2C", ascending=False).iloc[0]
    weakest = rows.sort_values("W2C", ascending=True).iloc[0]
    momentum = rows.sort_values("W2C Var", ascending=False).iloc[0]
    pressure = rows.sort_values("Rank", ascending=False).iloc[0]
    avg_w2c = rows["W2C"].mean()
    total_sales = rows["Sales"].sum()
    total_visitors = rows["Visitors"].sum()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — title
    slide = prs.slides.add_slide(blank)
    ppt_valtech_header(slide)
    ppt_add_textbox(slide, 0.72, 1.55, 11.4, 1.8, f"{selected_brand} Macro Conversion Storyteller Report", size=38, bold=False)
    ppt_add_textbox(
        slide, 0.76, 3.35, 10.8, 0.75,
        f"Website-to-contract conversion intelligence across UK, France, Germany, Italy and Spain. {CURRENT_LABEL} vs {PREVIOUS_LABEL}.",
        size=17, bold=False
    )
    rule = slide.shapes.add_shape(1, Inches(0.76), Inches(4.35), Inches(11.8), Inches(0.01))
    rule.fill.solid(); rule.fill.fore_color.rgb = RGBColor(217, 215, 209); rule.line.color.rgb = RGBColor(217, 215, 209)
    ppt_add_textbox(slide, 0.76, 4.55, 10.8, 0.4, "Sources: Marklines passenger car sales and Similarweb unique visitors", size=10, color=(111,111,111))

    # Slide 2 — executive takeaway
    slide = prs.slides.add_slide(blank)
    ppt_valtech_header(slide)
    ppt_add_textbox(slide, 0.72, 1.35, 11.8, 0.55, "Executive takeaway", size=30, bold=False)
    takeaway = (
        f"{selected_brand} is strongest in {best['Market']} at {best['W2C']:.2f}% W2C and weakest in {weakest['Market']} at {weakest['W2C']:.2f}% W2C. "
        f"The core commercial question is why conversion efficiency changes so materially across markets, especially where the brand has visitor scale but weaker relative ranking."
    )
    ppt_add_textbox(slide, 0.76, 2.05, 11.6, 0.85, takeaway, size=16, color=(0,0,0))
    ppt_add_valtech_card(slide, 0.76, 3.25, 2.25, 1.25, "Average W2C", f"{avg_w2c:.2f}%", "Five-market average")
    ppt_add_valtech_card(slide, 3.25, 3.25, 2.25, 1.25, "Sales", fmt_metric_number(total_sales), CURRENT_LABEL)
    ppt_add_valtech_card(slide, 5.74, 3.25, 2.25, 1.25, "Visitors", fmt_metric_number(total_visitors), CURRENT_LABEL)
    ppt_add_valtech_card(slide, 8.23, 3.25, 2.25, 1.25, "Strongest", str(best["Market"]), f"{best['W2C']:.2f}% W2C")
    ppt_add_valtech_card(slide, 10.72, 3.25, 2.25, 1.25, "Weakest", str(weakest["Market"]), f"{weakest['W2C']:.2f}% W2C", accent="red")

    # Slide 3 — five-market performance
    slide = prs.slides.add_slide(blank)
    ppt_valtech_header(slide)
    ppt_add_textbox(slide, 0.72, 1.25, 11.8, 0.55, "Five-market performance", size=30, bold=False)
    table_df = pd.DataFrame({
        "Market": rows["Market"],
        "W2C": rows["W2C"].map(lambda x: f"{x:.2f}%"),
        "Rank": rows.apply(lambda r: f"#{int(r['Rank'])} of {int(r['Peer Count'])}", axis=1),
        "Sales": rows["Sales"].map(fmt_metric_number),
        "Visitors": rows["Visitors"].map(fmt_metric_number),
        "Market gap": rows["Gap to Market Avg"].map(fmt_pp),
        "W2C movement": rows["W2C Var"].map(fmt_pp),
    })
    ppt_add_table(slide, table_df, 0.72, 2.1, 11.9, 2.3, font_size=9)
    ppt_add_textbox(slide, 0.76, 4.8, 11.4, 0.6, f"Rank scale is relative to the full OEM peer set in each market. Use this slide to separate absolute conversion strength from local competitive pressure.", size=14, color=(111,111,111))

    # Slide 4 — insight cards
    slide = prs.slides.add_slide(blank)
    ppt_valtech_header(slide)
    ppt_add_textbox(slide, 0.72, 1.25, 11.8, 0.55, "Storyteller insight cards", size=30, bold=False)
    market_gap = rows.sort_values("Gap to Market Avg").iloc[0]
    scale = rows.sort_values("Visitors", ascending=False).iloc[0]
    cards = [
        ("Strongest market", best["Market"], f"{best['W2C']:.2f}% W2C", "blue"),
        ("Weakest market", weakest["Market"], f"{weakest['W2C']:.2f}% W2C", "red"),
        ("Conversion gap", market_gap["Market"], fmt_pp(market_gap["Gap to Market Avg"]), "red"),
        ("Scale opportunity", scale["Market"], fmt_metric_number(scale["Visitors"]), "blue"),
        ("Competitive pressure", pressure["Market"], f"#{int(pressure['Rank'])} of {int(pressure['Peer Count'])}", "red"),
        ("Momentum", momentum["Market"], fmt_pp(momentum["W2C Var"]), "blue" if momentum["W2C Var"] >= 0 else "red"),
    ]
    xs = [0.76, 4.62, 8.48, 0.76, 4.62, 8.48]
    ys = [2.05, 2.05, 2.05, 4.05, 4.05, 4.05]
    for (title, market, value, accent), x, y in zip(cards, xs, ys):
        ppt_add_valtech_card(slide, x, y, 3.25, 1.35, title, value, market, accent=accent)

    # Slide 5 — action narrative
    slide = prs.slides.add_slide(blank)
    ppt_valtech_header(slide)
    ppt_add_textbox(slide, 0.72, 1.25, 11.8, 0.55, "Recommended diagnostic narrative", size=30, bold=False)
    bullets = [
        f"Use {best['Market']} as the reference market for strongest conversion behaviour and investigate what is transferable.",
        f"Prioritise {weakest['Market']} for conversion journey, offer, stock, retailer hand-off and demand-quality diagnosis.",
        f"Review {pressure['Market']} for competitive pressure because the selected OEM has the weakest relative rank there.",
        "Separate visitor scale from conversion quality: a high-audience market is not automatically a high-efficiency market.",
        "Treat W2C as a macro commercial-efficiency signal, not proof of website UX causality on its own.",
    ]
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.5), Inches(4.4))
    tf = box.text_frame; tf.clear()
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"• {ppt_safe_text(b)}"
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(0, 0, 0)
        para.space_after = Pt(10)

    # Slide 6 — methodology
    slide = prs.slides.add_slide(blank)
    ppt_valtech_header(slide)
    ppt_add_textbox(slide, 0.72, 1.25, 11.8, 0.55, "Methodology and caveats", size=30, bold=False)
    notes = [
        "Website-to-Contract Conversion Rate = passenger car sales / monthly deduplicated unique website visitors.",
        "Passenger car sales source: Marklines. Unique visitor source: Similarweb.",
        "The metric is designed for competitor intelligence and commercial diagnosis, not direct attribution.",
        "Use the report to identify where deeper investigation is needed across product, pricing, stock, retailer execution, media quality and website journey.",
    ]
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.5), Inches(4.4))
    tf = box.text_frame; tf.clear()
    for i, n in enumerate(notes):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"• {ppt_safe_text(n)}"
        para.font.size = Pt(17)
        para.font.color.rgb = RGBColor(0, 0, 0)
        para.space_after = Pt(10)

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()

def render_oem_insight_report_page(data):
    section("Insight Report")
    st.markdown(
        "<div class='report-builder-grid'>"
        "<div class='report-builder-card'>"
        "<div class='report-builder-title'>Step 1: Select the OEM you want to analyse</div>"
        "<div class='report-builder-copy'>Choose any OEM in the dataset. The same rules, metrics and ranking logic are applied to every brand.</div>"
        "</div>"
        "<div class='report-builder-card'>"
        "<div class='report-builder-title'>Step 2: Select the time period</div>"
        "<div class='report-builder-copy'>Use the comparison selector in the sidebar to switch between annual and Jan-Apr views, then review the generated five-market report below.</div>"
        "</div>"
        "</div>",
        unsafe_allow_html=True,
    )

    available = sorted(data["OEM"].dropna().astype(str).unique())
    default_idx = available.index("Honda") if "Honda" in available else 0
    selected_brand = st.selectbox("Step 1: Select OEM", available, index=default_idx, key="generic_insight_oem")
    st.info(f"Step 2: Active time period is {CURRENT_LABEL} vs {PREVIOUS_LABEL}. Change it in the sidebar under Comparison view.")

    rows = selected_oem_market_insight_rows(data, selected_brand)
    if rows.empty:
        st.warning("No insight data is available for the selected OEM and time period.")
        return

    try:
        ppt_bytes = create_oem_insight_storyteller_ppt(data, selected_brand, rows)
        st.download_button(
            "Download Storyteller PPT",
            data=ppt_bytes,
            file_name=f"{selected_brand.lower().replace(' ', '_')}_macro_conversion_storyteller_report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            key="download_oem_storyteller_ppt",
        )
    except Exception as exc:
        st.warning(f"PowerPoint export is unavailable: {exc}")

    best = rows.sort_values("W2C", ascending=False).iloc[0]
    weakest = rows.sort_values("W2C", ascending=True).iloc[0]
    momentum = rows.sort_values("W2C Var", ascending=False).iloc[0]
    pressure = rows.sort_values("Rank", ascending=False).iloc[0]
    avg_w2c = rows["W2C"].mean()
    total_sales = rows["Sales"].sum()
    total_visitors = rows["Visitors"].sum()

    st.markdown(
        "<div class='assistant-shell'>"
        f"<div class='assistant-title'>Executive takeaway — {selected_brand}</div>"
        f"<div class='assistant-copy'>{selected_brand} shows strongest Website-to-Contract performance in <b>{best['Market']}</b> at <b>{best['W2C']:.2f}%</b> and weakest performance in <b>{weakest['Market']}</b> at <b>{weakest['W2C']:.2f}%</b>. "
        f"The main opportunity is to close the conversion gap in {weakest['Market']}, while using {best['Market']} as the strongest cross-market reference point. "
        f"The strongest positive movement is in {momentum['Market']} at <b>{fmt_pp(momentum['W2C Var'])}</b> versus {PREVIOUS_LABEL}.</div>"
        "</div>",
        unsafe_allow_html=True,
    )

    kpi_html = "<div class='mm5-kpi-grid'>"
    kpis = [
        ("Average W2C rate", f"{avg_w2c:.2f}%", f"Across five markets"),
        ("Passenger car sales", fmt_metric_number(total_sales), CURRENT_LABEL),
        ("Unique visitors", fmt_metric_number(total_visitors), CURRENT_LABEL),
        ("Strongest market", best["Market"], f"{best['W2C']:.2f}% W2C"),
        ("Weakest market", weakest["Market"], f"{weakest['W2C']:.2f}% W2C"),
    ]
    for i, (label, value, delta) in enumerate(kpis):
        cls = "mm5-kpi-card primary" if i == 0 else "mm5-kpi-card"
        kpi_html += f"<div class='{cls}'><div class='mm5-kpi-label'>{label}</div><div class='mm5-kpi-value'>{value}</div><div style='margin-top:10px;color:#6F7782;font-weight:800;'>{delta}</div></div>"
    kpi_html += "</div>"
    st.markdown(kpi_html, unsafe_allow_html=True)

    section("Market performance cards")
    card_html = "<div class='brand-insight-card-grid'>"
    for _, r in rows.iterrows():
        card_html += (
            "<div class='brand-insight-card'>"
            f"<div class='brand-insight-title'>{market_label_with_flag(r['Market'])}</div>"
            "<div class='brand-insight-label'>W2C rate</div>"
            f"<div class='brand-insight-main'>{r['W2C']:.2f}%</div>"
            f"<div class='brand-insight-meta'>Rank: <b>{int(r['Rank'])} of {int(r['Peer Count'])}</b><br>"
            f"Passenger car sales: <b>{fmt_metric_number(r['Sales'])}</b> {badge_span(r['Sales YoY'])}<br>"
            f"Unique visitors: <b>{fmt_metric_number(r['Visitors'])}</b> {badge_span(r['Visitors YoY'])}<br>"
            f"Market-average gap: {badge_span(r['Gap to Market Avg'], suffix='pp')}</div>"
            "</div>"
        )
    card_html += "</div>"
    st.markdown(card_html, unsafe_allow_html=True)

    section("Cross-market visualisations")
    c1, c2 = st.columns(2)
    with c1:
        fig = go.Figure(go.Bar(x=rows["Market"], y=rows["W2C"], text=rows["W2C"].map(lambda x: f"{x:.2f}%"), textposition="outside", marker_color=VALTECH_BLUE))
        fig.update_layout(title="Website-to-Contract Conversion Rate by market", yaxis_title="W2C rate", xaxis_title="Market", height=430, plot_bgcolor=WHITE, paper_bgcolor=WHITE)
        fig.update_yaxes(ticksuffix="%", gridcolor="#EDF1F5")
        st.plotly_chart(fig, use_container_width=True)
    with c2:
        fig = go.Figure(go.Bar(x=rows["Market"], y=rows["Rank"], text=rows["Rank"].map(lambda x: f"#{int(x)}"), textposition="outside", marker_color=VALTECH_GREY))
        fig.update_layout(title="Market ranking view", yaxis_title="Rank within market", xaxis_title="Market", height=430, plot_bgcolor=WHITE, paper_bgcolor=WHITE)
        max_rank_tick = int(math.ceil(rows["Rank"].max() / 5) * 5)
        fig.update_yaxes(autorange=False, range=[max_rank_tick, 0], dtick=5, tick0=0, gridcolor="#EDF1F5")
        st.plotly_chart(fig, use_container_width=True)

    c3, c4 = st.columns(2)
    with c3:
        fig = go.Figure(go.Bar(x=rows["Market"], y=rows["W2C Var"], text=rows["W2C Var"].map(fmt_pp), textposition="outside", marker_color=rows["W2C Var"].map(lambda x: GREEN if x >= 0 else PINK)))
        fig.update_layout(title=f"W2C movement vs {PREVIOUS_LABEL}", yaxis_title="Percentage-point movement", xaxis_title="Market", height=430, plot_bgcolor=WHITE, paper_bgcolor=WHITE)
        fig.update_yaxes(ticksuffix="pp", gridcolor="#EDF1F5", zeroline=True, zerolinecolor="#B8C0C8")
        st.plotly_chart(fig, use_container_width=True)
    with c4:
        fig = go.Figure(go.Scatter(
            x=rows["Visitors"],
            y=rows["W2C"],
            mode="markers+text",
            text=rows["Market"],
            textposition="top center",
            marker=dict(size=rows["Sales"].map(lambda x: max(18, math.sqrt(x / max(rows["Sales"].max(), 1)) * 70)), color=VALTECH_BLUE, opacity=0.78, line=dict(width=2, color="white")),
            hovertemplate="<b>%{text}</b><br>Visitors: %{x:,.0f}<br>W2C: %{y:.2f}%<extra></extra>",
        ))
        fig.update_layout(title="Opportunity matrix: visitor scale vs conversion", xaxis_title="Unique visitors", yaxis_title="W2C rate", height=430, plot_bgcolor=WHITE, paper_bgcolor=WHITE)
        fig.update_xaxes(tickformat=",", gridcolor="#EDF1F5")
        fig.update_yaxes(ticksuffix="%", gridcolor="#EDF1F5")
        st.plotly_chart(fig, use_container_width=True)

    section("Generated insight cards")
    cards = [
        assistant_card("Strongest market", f"{best['Market']} has the highest W2C rate for {selected_brand} across the five markets.", f"{best['W2C']:.2f}%", best["Market"], "positive"),
        assistant_card("Weakest market", f"{weakest['Market']} is the lowest-converting market and should be prioritised for funnel diagnosis.", f"{weakest['W2C']:.2f}%", weakest["Market"], "risk"),
        assistant_card("Conversion opportunity", f"The largest negative market-average gap is in {rows.sort_values('Gap to Market Avg').iloc[0]['Market']}.", fmt_pp(rows.sort_values('Gap to Market Avg').iloc[0]['Gap to Market Avg']), "benchmark gap", "risk"),
        assistant_card("Scale opportunity", f"{rows.sort_values('Visitors', ascending=False).iloc[0]['Market']} has the largest visitor base for {selected_brand}.", fmt_metric_number(rows.sort_values('Visitors', ascending=False).iloc[0]['Visitors']), "visitor scale", "primary"),
        assistant_card("Competitive pressure", f"Weakest relative rank is in {pressure['Market']}. Review the local peer set and leader behaviour there.", f"#{int(pressure['Rank'])} of {int(pressure['Peer Count'])}", pressure["Market"], "neutral"),
        assistant_card("Momentum", f"{momentum['Market']} has the strongest positive W2C movement versus the prior period.", fmt_pp(momentum["W2C Var"]), f"vs {PREVIOUS_LABEL}", "positive" if momentum["W2C Var"] >= 0 else "risk"),
    ]
    st.markdown(assistant_result_html(cards), unsafe_allow_html=True)

    table = rows.copy()
    table["W2C"] = table["W2C"].map(lambda x: f"{x:.2f}%")
    table["W2C Var"] = table["W2C Var"].map(fmt_pp)
    table["Sales"] = table["Sales"].map(fmt_metric_number)
    table["Visitors"] = table["Visitors"].map(fmt_metric_number)
    table["Gap to Leader"] = table["Gap to Leader"].map(fmt_pp)
    st.download_button(
        "Download selected OEM insight data as CSV",
        data=rows.to_csv(index=False).encode("utf-8"),
        file_name=f"{selected_brand.lower().replace(' ', '_')}_insight_report.csv",
        mime="text/csv",
    )
    render_footer()


def render_insight_report_page(data):
    render_oem_insight_report_page(data)


def render_data_assistant_page(data, selected_oems):
    render_data_assistant(data, selected_oems)
    render_footer()




def render_start_here_page(data):
    section("Start here — how to use this dashboard")

    included_markets = ["UK", "France", "Germany", "Italy", "Spain"]
    period_labels = [cfg["coverage_label"] for cfg in COMPARISON_OPTIONS.values()]
    category_labels = available_oem_categories()
    oems = sorted(data["OEM"].dropna().astype(str).unique().tolist())

    intro_html = (
        "<div class='start-intro'>"
        "<b>This is the OEM Macro Conversion Intelligence Report.</b><br><br>"
        "Use this dashboard to compare how OEM website demand converts into passenger car sales across the five included markets. "
        "The core metric is <b>Website-to-Contract Conversion Rate</b>: passenger car sales divided by monthly deduplicated unique website visitors.<br><br>"
        "The dashboard should be read as a macro commercial-efficiency tool. It helps identify where an OEM may be winning through demand quality, product-market fit, offer strength, retailer execution, stock availability, or conversion journey performance."
        "</div>"
    )
    st.markdown(intro_html, unsafe_allow_html=True)

    section("Report coverage")
    coverage_html = (
        "<div class='methodology-grid'>"
        "<div class='methodology-card'>"
        "<div class='methodology-callout'>Markets included</div>"
        "<h4>Five-market view</h4>"
        f"<div class='methodology-formula'><b>Markets:</b><br>{', '.join(included_markets)}</div>"
        "<ul><li>MM5 is available as an aggregated view.</li><li>Individual market pages show UK, France, Germany, Italy and Spain.</li></ul>"
        "</div>"
        "<div class='methodology-card'>"
        "<div class='methodology-callout'>Time periods included</div>"
        "<h4>Available comparisons</h4>"
        f"<div class='methodology-formula'><b>Comparison views:</b><br>{' · '.join(COMPARISON_OPTIONS.keys())}</div>"
        f"<ul><li>Coverage labels: {' · '.join(period_labels)}</li><li>Use the sidebar to switch comparison period.</li></ul>"
        "</div>"
        "<div class='methodology-card'>"
        "<div class='methodology-callout'>OEM groupings included</div>"
        "<h4>Competitor categories</h4>"
        f"<div class='methodology-formula'><b>Groupings:</b><br>{' · '.join(category_labels)}</div>"
        "<ul><li>Categories support like-for-like competitor analysis.</li><li>No OEM is locked, pinned, or given special treatment.</li></ul>"
        "</div>"
        "<div class='methodology-card'>"
        "<div class='methodology-callout'>Primary metric</div>"
        "<h4>Website-to-Contract Conversion Rate</h4>"
        "<div class='methodology-formula'><b>Formula</b><br>Passenger car sales / monthly deduplicated unique website visitors.</div>"
        "<ul><li>Sales source: Marklines.</li><li>Unique visitor source: Similarweb.</li><li>Interpret as commercial efficiency, not purely website UX.</li></ul>"
        "</div>"
        "</div>"
    )
    st.markdown(coverage_html, unsafe_allow_html=True)

    section("Factors that can influence Website-to-Contract Conversion Rate")
    factors = [
        "Brand Strategy",
        "Marketing Strategy",
        "Ecommerce Strategy",
        "Website Usability",
        "Website Experience",
        "Pricing & Promotions Strategy",
        "Ownership Options",
        "Product Strategy",
        "Stock Availability",
        "Network Strategy",
        "Powertrain Strategy",
        "Macro Economic factors",
    ]
    cards_html = "".join([f"<div class='start-factor-card'>{factor}</div>" for factor in factors])
    st.markdown(f"<div class='start-factors-grid'>{cards_html}</div>", unsafe_allow_html=True)

    section("How to interpret Similarweb unique visitors")
    c1, c2 = st.columns(2)
    with c1:
        method_html = (
            "<div class='method-card'>"
            "<div class='method-title'>Why Similarweb will not exactly match web analytics</div>"
            "<div class='method-copy'>"
            "Similarweb traffic estimates are built from a blended methodology and should be treated as a consistent external market benchmark, "
            "not as a direct replacement for first-party analytics. The direction and relative ranking are usually more useful than trying to reconcile every absolute number."
            "</div>"
            "</div>"
        )
        st.markdown(method_html, unsafe_allow_html=True)
    with c2:
        consent_html = (
            "<div class='method-card'>"
            "<div class='method-title'>Why consent management affects sources differently</div>"
            "<div class='method-copy'>"
            "First-party analytics and external market-estimation tools are affected by different collection methods. Use the dashboard as a competitor intelligence lens, "
            "then validate priority hypotheses with owned analytics, CRM, lead handling and retail performance data."
            "</div>"
            "</div>"
        )
        st.markdown(consent_html, unsafe_allow_html=True)

    section("OEMs included in this report")
    pills = "".join([f"<span class='oem-pill'>{o}</span>" for o in oems])
    st.markdown(f"<div class='oem-pill-wrap'>{pills}</div>", unsafe_allow_html=True)

    render_footer()

def usecase_card(audience, title, copy, reports, audience_class="shared"):
    reports_html = "<br>".join([f"• {r}" for r in reports])
    return (
        f"<div class='usecase-card {audience_class}'>"
        f"<div class='usecase-audience'>{audience}</div>"
        f"<div class='usecase-title'>{title}</div>"
        f"<div class='usecase-copy'>{copy}</div>"
        f"<div class='usecase-report'><b>Recommended report view</b><br>{reports_html}</div>"
        f"</div>"
    )


def render_use_cases_page(data):
    section("Use cases — what to use this dashboard for")

    st.markdown(
        "<div class='usecase-note'>"
        "Use this report to compare OEM performance across markets, identify conversion opportunities, understand where brands are gaining or losing competitive advantage, and frame better planning questions for product, marketing, ecommerce, network and sales teams."
        "</div>",
        unsafe_allow_html=True,
    )

    use_cases = [
        ("Prioritise markets", "Identify where a selected OEM is strongest, weakest, and most exposed across the five included markets."),
        ("Compare competitor groups", "Use OEM categories to compare volume leaders, premium brands, EV challengers, mainstream brands and new entrants consistently."),
        ("Diagnose traffic quality", "Separate visitor growth from sales conversion to spot where demand is not translating into contracts."),
        ("Find conversion opportunities", "Use W2C gaps and movement to challenge journey friction, stock visibility, offers, retailer hand-off and lead quality."),
        ("Prepare executive updates", "Turn cross-market signals into a concise story about risk, opportunity, strongest markets and weakest markets."),
        ("Shape quarterly planning", "Use market and scorecard outputs to decide where deeper investigation is needed before commercial decisions."),
    ]

    html = "<div class='flow-usecase-grid'>"
    for i, (title, copy) in enumerate(use_cases, start=1):
        html += (
            "<div class='flow-usecase-card'>"
            f"<div class='flow-number'>{i}</div>"
            f"<div class='flow-title'>{title}</div>"
            f"<div class='flow-copy'>{copy}</div>"
            "</div>"
        )
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)

    section("Deeper insights into website strategies")
    links = [
        ("🇬🇧", "UK", "https://valtech-uk-auto.netlify.app/"),
        ("🇩🇪", "Germany", "https://valtech-de-auto.netlify.app/"),
        ("🇫🇷", "France", "https://valtech-fr-auto.netlify.app/"),
        ("🇮🇹", "Italy", "https://valtech-it-auto.netlify.app/"),
        ("🇪🇸", "Spain", "https://valtech-es-auto.netlify.app/"),
    ]
    html = "<div class='dashboard-link-grid'>"
    for flag, name, url in links:
        html += f"<a class='dashboard-link-card' href='{url}' target='_blank'>{flag} {name}<span>Open market website strategy dashboard</span></a>"
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)

    render_footer()

def render_gap_analysis_page(data, market):
    section(f"Gap Analysis — {market}")
    st.caption(f"Comparison selected: {CURRENT_LABEL} vs {PREVIOUS_LABEL}.")
    st.caption(f"All headline metrics reflect current period performance; variance badges show movement versus previous period.")
    render_brand_strip()

    combined = market_year(data, market, 2025, ["Toyota", "Lexus"])
    combined_24 = market_year(data, market, 2024, ["Toyota", "Lexus"])

    if combined.empty:
        st.warning("No selected OEM data available.")
        return

    sales_25 = combined["Sales"].sum()
    visitors_25 = combined["UniqueVisitors"].sum()
    conv_25 = sales_25 / visitors_25 * 100
    sales_24 = combined_24["Sales"].sum()
    visitors_24 = combined_24["UniqueVisitors"].sum()
    conv_24 = sales_24 / visitors_24 * 100 if visitors_24 else 0

    c1, c2, c3, c4 = st.columns(4)
    c1.metric(f"selected OEM Passenger Car Sales", fmt_int(sales_25), f"{fmt_pct((sales_25 / sales_24 - 1) * 100)} vs previous period" if sales_24 else "n/a")
    c2.metric(f"selected OEM Unique Visitors", fmt_int(visitors_25), f"{fmt_pct((visitors_25 / visitors_24 - 1) * 100)} vs previous period" if visitors_24 else "n/a")
    c3.metric(f"Website-to-Contract Conversion Rate", f"{conv_25:.2f}%", f"{fmt_pp(conv_25 - conv_24)} vs previous period")
    c4.metric(f"Visits per contract", fmt_int(visitors_25 / sales_25))

    render_brand_detail(data, market)
    render_market_weakness(data)
    render_toyota_lexus_recommendations(data)
    render_footer()




def signed_class(value):
    """Return positive/negative/neutral from a numeric value or formatted delta string."""
    try:
        v = float(value)
    except Exception:
        text_value = str(value).strip()
        if text_value.startswith('+'):
            return 'positive'
        if text_value.startswith('-'):
            return 'negative'
        return 'neutral'
    if v > 0:
        return 'positive'
    if v < 0:
        return 'negative'
    return 'neutral'


def styled_delta_html(value, label=None, suffix=''):
    cls = signed_class(value)
    if label is None:
        label = str(value)
    return f"<span class='{cls}'>{label}{suffix}</span>"


def exec_kpi_card(label, value, delta=None, delta_class=None):
    if delta is not None:
        cls = delta_class or signed_class(delta)
        delta_html = f"<div class='exec-kpi-delta {cls}'>{delta}</div>"
    else:
        delta_html = ""
    return (
        f"<div class='exec-kpi-card'>"
        f"<div class='exec-kpi-label'>{label}</div>"
        f"<div class='exec-kpi-value'>{value}</div>"
        f"{delta_html}"
        f"</div>"
    )



def posneg_class_text(value_text):
    text_value = str(value_text)
    if text_value.startswith("+"):
        return "pos"
    if text_value.startswith("-"):
        return "neg"
    return "neutral"


def badge_span(value, suffix="%"):
    if value is None or pd.isna(value):
        text_value = "n/a"
    elif suffix == "pp":
        text_value = fmt_pp(value)
    elif suffix == "raw":
        text_value = str(value)
    else:
        text_value = fmt_pct(value)
    cls = posneg_class_text(text_value)
    bg = "#DDF8EC" if cls == "pos" else "#FFE5EF" if cls == "neg" else "#EEF2F6"
    col = "#12C76B" if cls == "pos" else "#FF2F6D" if cls == "neg" else "#6F6F6F"
    return f"<span style='background:{bg};color:{col};padding:5px 9px;border-radius:999px;font-weight:850;'>{text_value}</span>"


def html_table(df, highlight_brands=None, sticky_brand=True):
    highlight_brands = set(highlight_brands or [])
    html = "<div class='html-table-wrap'><table class='html-table'><thead><tr>"
    for col in df.columns:
        cls = " class='sticky-brand'" if sticky_brand and col == "Brand" else ""
        html += f"<th{cls}>{col}</th>"
    html += "</tr></thead><tbody>"
    for _, r in df.iterrows():
        brand_text = str(r.get("Brand", "")).replace("★ ", "")
        tr_cls = " class='table-highlight'" if brand_text in highlight_brands else ""
        html += f"<tr{tr_cls}>"
        for col in df.columns:
            cls = " class='sticky-brand'" if sticky_brand and col == "Brand" else ""
            html += f"<td{cls}>{r[col]}</td>"
        html += "</tr>"
    html += "</tbody></table></div>"
    st.markdown(html, unsafe_allow_html=True)


def medal_for_rank(rank):
    try:
        rank = int(rank)
    except Exception:
        return str(rank)
    return {1: "🥇 1", 2: "🥈 2", 3: "🥉 3"}.get(rank, str(rank))


def market_kpi_summary(data, market, selected_oems=None):
    current = market_year(data, market, 2025, selected_oems)
    previous = market_year(data, market, 2024, selected_oems)
    if current.empty:
        return None
    sales_25 = current["Sales"].sum()
    visitors_25 = current["UniqueVisitors"].sum()
    conv_25 = sales_25 / visitors_25 * 100 if visitors_25 else 0
    visits_sale = visitors_25 / sales_25 if sales_25 else None
    sales_24 = previous["Sales"].sum() if not previous.empty else 0
    visitors_24 = previous["UniqueVisitors"].sum() if not previous.empty else 0
    conv_24 = sales_24 / visitors_24 * 100 if visitors_24 else 0
    return {
        "sales": sales_25,
        "visitors": visitors_25,
        "conv": conv_25,
        "visits_sale": visits_sale,
        "sales_delta": (sales_25 / sales_24 - 1) * 100 if sales_24 else None,
        "visitor_delta": (visitors_25 / visitors_24 - 1) * 100 if visitors_24 else None,
        "conv_delta": conv_25 - conv_24,
    }



MARKET_FLAGS = {
    "MM5": "🇪🇺",
    "UK": "🇬🇧",
    "France": "🇫🇷",
    "Germany": "🇩🇪",
    "Italy": "🇮🇹",
    "Spain": "🇪🇸",
}


def market_label_with_flag(market):
    return f"<span class='flag'>{MARKET_FLAGS.get(market, '')}</span>{market}"


def brand_competitor_set(brand, all_oems=None):
    if brand == "Lexus":
        base = list(LEXUS_SET)
    elif brand == "Toyota":
        base = list(TOYOTA_SET)
    else:
        base = []
    if brand not in base:
        base.append(brand)
    if all_oems is not None:
        base = [x for x in base if x in all_oems]
    return base


def brand_market_cards(data, brand):
    markets = ["UK", "France", "Germany", "Italy", "Spain"]
    html = "<div class='brand-insight-card-grid'>"
    for market in markets:
        yoy = yoy_table(data, market, None)
        row = get_row(yoy, brand)
        if row is None:
            continue
        html += (
            "<div class='brand-insight-card'>"
            f"<div class='brand-insight-title'>{market_label_with_flag(market)} {brand}</div>"
            "<div class='brand-insight-label'>W2C rate</div>"
            f"<div class='brand-insight-main'>{row['ConversionPct_2025']:.2f}%</div>"
            f"<div class='brand-insight-meta'>W2C movement: {badge_span(row['Conv Var pp'], suffix='pp')} vs {PREVIOUS_LABEL}</div>"
            f"<div class='brand-insight-meta' style='margin-top:10px;'>Passenger Car Sales: <b>{fmt_metric_number(row['Sales_2025'])}</b> {badge_span(row['Sales YoY %'])}<br>"
            f"Unique visitors: <b>{fmt_metric_number(row['UniqueVisitors_2025'])}</b> {badge_span(row['Visitors YoY %'])}</div>"
            "</div>"
        )
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)


def brand_gap_table(data, brand):
    rows = []
    competitor_set = brand_competitor_set(brand, sorted(data["OEM"].unique()))
    for market in ["UK", "France", "Germany", "Italy", "Spain"]:
        yoy = yoy_table(data, market, None)
        row = get_row(yoy, brand)
        cohort = yoy[yoy["OEM"].isin(competitor_set)].copy()
        if row is None or cohort.empty:
            continue
        leader = cohort.sort_values("ConversionPct_2025", ascending=False).iloc[0]
        rows.append({
            "Market": market_label_with_flag(market),
            "Brand": brand,
            "Visitor YoY": badge_span(row["Visitors YoY %"]),
            "Sales YoY": badge_span(row["Sales YoY %"]),
            "W2C rate": f"{row['ConversionPct_2025']:.2f}%",
            "Benchmark leader": leader["OEM"],
            "Leader W2C rate": f"{leader['ConversionPct_2025']:.2f}%",
            "Gap to leader": badge_span(row["ConversionPct_2025"] - leader["ConversionPct_2025"], suffix="pp"),
        })
    if rows:
        html_table(pd.DataFrame(rows), highlight_brands=[brand], sticky_brand=False)
    else:
        st.info(f"No {brand} insight data available.")


def render_brand_insights_bubble(data, brand):
    competitor_set = brand_competitor_set(brand, sorted(data["OEM"].unique()))
    x_axis_metric = st.selectbox("Bubble chart x axis", ["Unique Visitors", "Passenger Car Sales"], index=0, key=f"{brand.lower()}_insight_bubble_x")
    size_label = "Passenger Car Sales" if x_axis_metric == "Unique Visitors" else "Unique visitors"
    st.markdown(f"<div class='bubble-key'>Bubble size represents <b>{size_label}</b>. Each market view is pre-set to the relevant selected competitor set.</div>", unsafe_allow_html=True)
    tabs = st.tabs(["UK", "France", "Germany", "Italy", "Spain"])
    for market, tab in zip(["UK", "France", "Germany", "Italy", "Spain"], tabs):
        with tab:
            df = data[(data["Market"] == market) & (data["OEM"].isin(competitor_set))].copy()
            if df.empty:
                st.info("No bubble chart data for this market.")
                continue
            fig = build_bubble_chart(df, competitor_set, market, "Previous and current + shift", False, x_axis_metric=x_axis_metric, view_label="OEM")
            st.plotly_chart(fig, use_container_width=True)


def render_brand_insights_page(data, brand):
    section(f"{brand} Insights")
    render_brand_logo_header(brand)
    st.caption(f"Pre-set to the relevant selected competitor set. Comparison selected: {CURRENT_LABEL} vs {PREVIOUS_LABEL}.")
    section(f"{brand} market performance")
    brand_market_cards(data, brand)
    section(f"{brand} competitor-set market benchmark")
    competitor_market_cards(data, brand)
    section(f"{brand} market gap analysis")
    brand_gap_table(data, brand)
    section(f"{brand} competitor bubble chart by market")
    render_brand_insights_bubble(data, brand)
    render_footer()



def oems_for_cluster(data, cluster_name):
    return sorted([o for o in data["OEM"].unique() if cluster_for_oem(o) == cluster_name])


def resolve_tme_preset(data, preset_name):
    if preset_name == "Toyota volume competitors":
        return [o for o in TOYOTA_SET if o in set(data["OEM"].unique())]
    if preset_name == "Lexus premium competitors":
        return [o for o in LEXUS_SET if o in set(data["OEM"].unique())]
    return sorted(data["OEM"].unique())


def default_volume_oems(data):
    volume = oems_for_cluster(data, "Volume Leaders")
    return volume if volume else sorted(data["OEM"].unique())[:10]

def render_brand_logo_header(brand):
    logo = LEXUS_LOGO if brand == "Lexus" else TOYOTA_LOGO
    height = "56px" if brand == "Lexus" else "66px"
    st.markdown(
        f"<div class='brand-logo-header'><img src='{logo}' style='height:{height};' alt='{brand} logo'></div>",
        unsafe_allow_html=True,
    )


def render_mm5_summary_kpi_cards(data, selected_oems):
    summary = market_kpi_summary(data, "MM5", selected_oems)
    if not summary:
        st.warning("No MM5 data available for this selection.")
        return False

    cards = [
        ("MM5 Passenger Car Sales", fmt_metric_number(summary["sales"]), badge_span(summary["sales_delta"]) if summary["sales_delta"] is not None else "n/a"),
        ("MM5 unique visitors", fmt_metric_number(summary["visitors"]), badge_span(summary["visitor_delta"]) if summary["visitor_delta"] is not None else "n/a"),
        ("MM5 W2C rate", f"{summary['conv']:.2f}%", badge_span(summary["conv_delta"], suffix="pp")),
        ("MM5 visits per sale", fmt_int(summary["visits_sale"]) if summary["visits_sale"] else "n/a", "lower is better"),
        ("Coverage", CURRENT_LABEL, f"vs {PREVIOUS_LABEL}"),
    ]
    html = "<div class='mm5-kpi-grid'>"
    for i, (label, value, delta) in enumerate(cards):
        cls = "mm5-kpi-card primary" if i == 0 else "mm5-kpi-card"
        html += (
            f"<div class='{cls}'>"
            f"<div class='mm5-kpi-label'>{label}</div>"
            f"<div class='mm5-kpi-value'>{value}</div>"
            f"<div style='margin-top:10px;'>{delta}</div>"
            "</div>"
        )
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)
    return True


def competitor_market_cards(data, brand):
    competitor_set = brand_competitor_set(brand, sorted(data["OEM"].unique()))
    html = "<div class='brand-insight-card-grid'>"
    for market in ["UK", "France", "Germany", "Italy", "Spain"]:
        summary = market_kpi_summary(data, market, competitor_set)
        if not summary:
            continue
        html += (
            "<div class='brand-insight-card competitor'>"
            f"<div class='brand-insight-title'>{market_label_with_flag(market)} competitor set</div>"
            "<div class='brand-insight-label'>Benchmark W2C rate</div>"
            f"<div class='brand-insight-main'>{summary['conv']:.2f}%</div>"
            f"<div class='brand-insight-meta'>Passenger Car Sales: <b>{fmt_metric_number(summary['sales'])}</b><br>"
            f"Unique visitors: <b>{fmt_metric_number(summary['visitors'])}</b><br>"
            f"Sales movement: {badge_span(summary['sales_delta']) if summary['sales_delta'] is not None else 'n/a'}</div>"
            "</div>"
        )
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)


def sort_scorecard(score, sort_by, descending):
    if sort_by == "W2C rate":
        col = "ConversionPct_2025"
    elif sort_by == "W2C variance":
        col = "Conv Var pp"
    elif sort_by == "Passenger Car Sales":
        col = "Sales_2025"
    elif sort_by == "Sales YoY":
        col = "Sales YoY %"
    elif sort_by == "Unique visitors":
        col = "UniqueVisitors_2025"
    elif sort_by == "Visitor YoY":
        col = "Visitors YoY %"
    elif sort_by == "Visits to sale":
        col = "Visits to Sale 2025"
    else:
        col = "W2C Ranking"
        descending = False
    return score.sort_values(col, ascending=not descending)

def render_market_kpi_cards(data, market, selected_oems):
    if market != "MM5":
        return render_exec_kpis(data, market, selected_oems)

    if not render_mm5_summary_kpi_cards(data, selected_oems):
        return False

    html = "<div class='market-five-grid'>"
    for m in ["UK", "France", "Germany", "Italy", "Spain"]:
        summary = market_kpi_summary(data, m, selected_oems)
        if not summary:
            continue
        html += (
            "<div class='exec-market-card'>"
            f"<div class='exec-market-title'>{market_label_with_flag(m)}</div>"
            f"<div class='exec-market-row'><span class='exec-market-label'>Passenger Car Sales</span><span class='exec-market-value'>{fmt_metric_number(summary['sales'])}</span></div>"
            f"<div class='exec-market-row'><span class='exec-market-label'>Unique visitors</span><span class='exec-market-value'>{fmt_metric_number(summary['visitors'])}</span></div>"
            f"<div class='exec-market-row'><span class='exec-market-label'>Market W2C rate</span><span class='exec-market-value'>{summary['conv']:.2f}%</span></div>"
            f"<div class='exec-market-row'><span class='exec-market-label'>Visits per sale</span><span class='exec-market-value'>{fmt_int(summary['visits_sale']) if summary['visits_sale'] else 'n/a'}</span></div>"
            f"<div class='exec-market-row'><span class='exec-market-label'>Sales movement</span><span class='exec-market-value'>{badge_span(summary['sales_delta']) if summary['sales_delta'] is not None else 'n/a'}</span></div>"
            "</div>"
        )
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)
    return True


def render_scorecard_podiums(score):
    top_sales = score.sort_values("Sales_2025", ascending=False).head(3)
    top_eff = score.sort_values("ConversionPct_2025", ascending=False).head(3)
    def podium_html(title, df, value_col, formatter):
        html = f"<div class='score-podium-card'><div class='score-podium-title'>{title}</div>"
        for i, (_, r) in enumerate(df.iterrows(), start=1):
            html += (
                "<div class='score-podium-row'>"
                f"<div>{['🥇','🥈','🥉'][i-1]}</div>"
                f"<div class='score-podium-brand'>{r['OEM']}</div>"
                f"<div class='score-podium-value'>{formatter(r[value_col])}</div>"
                "</div>"
            )
        html += "</div>"
        return html
    html = "<div class='score-top-grid'>" + podium_html("Top 3 Sales Volume", top_sales, "Sales_2025", fmt_metric_number) + podium_html("Top 3 Sales Efficiency", top_eff, "ConversionPct_2025", lambda x: f"{x:.2f}%") + "</div>"
    st.markdown(html, unsafe_allow_html=True)

def render_exec_kpis(data, market, selected_oems):
    summary = market_kpi_summary(data, market, selected_oems)
    if not summary:
        st.warning("No data available for this selection.")
        return False
    cards = [
        exec_kpi_card(f"{CURRENT_LABEL} Passenger Car Sales", fmt_metric_number(summary["sales"]), badge_span(summary["sales_delta"]) if summary["sales_delta"] is not None else "n/a"),
        exec_kpi_card(f"{CURRENT_LABEL} unique visitors", fmt_metric_number(summary["visitors"]), badge_span(summary["visitor_delta"]) if summary["visitor_delta"] is not None else "n/a"),
        exec_kpi_card("Market W2C rate", f"{summary['conv']:.2f}%", badge_span(summary["conv_delta"], suffix="pp")),
        exec_kpi_card("Visits per sale", fmt_int(summary["visits_sale"]) if summary["visits_sale"] else "n/a", "lower is better"),
    ]
    st.markdown("<div class='exec-kpi-grid'>" + "".join(cards) + "</div>", unsafe_allow_html=True)
    return True


def render_toyota_lexus_market_share(data, market):
    current_all = market_year(data, market, 2025, None)
    previous_all = market_year(data, market, 2024, None)
    current_tl = market_year(data, market, 2025, ["Toyota", "Lexus"])
    previous_tl = market_year(data, market, 2024, ["Toyota", "Lexus"])

    if current_all.empty or current_tl.empty:
        st.info("No selected OEM share data available for this market.")
        return

    all_visitors = current_all["UniqueVisitors"].sum()
    all_sales = current_all["Sales"].sum()
    tl_visitors = current_tl["UniqueVisitors"].sum()
    tl_sales = current_tl["Sales"].sum()

    prev_all_visitors = previous_all["UniqueVisitors"].sum() if not previous_all.empty else 0
    prev_all_sales = previous_all["Sales"].sum() if not previous_all.empty else 0
    prev_tl_visitors = previous_tl["UniqueVisitors"].sum() if not previous_tl.empty else 0
    prev_tl_sales = previous_tl["Sales"].sum() if not previous_tl.empty else 0

    visitor_share = tl_visitors / all_visitors * 100 if all_visitors else 0
    contract_share = tl_sales / all_sales * 100 if all_sales else 0
    prev_visitor_share = prev_tl_visitors / prev_all_visitors * 100 if prev_all_visitors else 0
    prev_contract_share = prev_tl_sales / prev_all_sales * 100 if prev_all_sales else 0

    visitor_share_delta = visitor_share - prev_visitor_share
    contract_share_delta = contract_share - prev_contract_share
    html = (
        "<div class='share-grid'>"
        "<div class='share-card'>"
        "<div class='share-title'>selected OEM share of unique visitors</div>"
        f"<div class='share-metric-row'><span class='share-label'>Share of market visitors</span><span class='share-value'>{visitor_share:.1f}%</span></div>"
        f"<div class='share-metric-row'><span class='share-label'>Unique visitors</span><span class='share-value'>{fmt_metric_number(tl_visitors)}</span></div>"
        f"<div class='share-metric-row'><span class='share-label'>Movement vs {PREVIOUS_LABEL}</span><span class='share-value {signed_class(visitor_share_delta)}'>{fmt_pp(visitor_share_delta)}</span></div>"
        "</div>"
        "<div class='share-card'>"
        "<div class='share-title'>selected OEM share of Passenger Car Sales</div>"
        f"<div class='share-metric-row'><span class='share-label'>Share of Passenger Car Sales</span><span class='share-value'>{contract_share:.1f}%</span></div>"
        f"<div class='share-metric-row'><span class='share-label'>Customer contracts</span><span class='share-value'>{fmt_metric_number(tl_sales)}</span></div>"
        f"<div class='share-metric-row'><span class='share-label'>Movement vs {PREVIOUS_LABEL}</span><span class='share-value {signed_class(contract_share_delta)}'>{fmt_pp(contract_share_delta)}</span></div>"
        "</div>"
        "</div>"
    )
    st.markdown(html, unsafe_allow_html=True)


def top10_leaders_table(data, market, selected_oems):
    yoy = yoy_table(data, market, selected_oems)
    if yoy.empty:
        st.info("No Top 10 data available for this selection.")
        return
    base = yoy.sort_values("ConversionPct_2025", ascending=False).reset_index(drop=True)
    base["Rank"] = range(1, len(base) + 1)
    table = base.head(10).copy()
    visitor_col = f"Visitor change vs {PREVIOUS_LABEL}"
    sales_col = f"Sales change vs {PREVIOUS_LABEL}"
    display = pd.DataFrame({
        "Rank": table["Rank"].map(medal_for_rank),
        "Brand": table["OEM"],
        "Category": table["OEM"].map(cluster_for_oem),
        "W2C Conv Rate": table["ConversionPct_2025"].map(lambda x: f"{x:.2f}%"),
        "Passenger Car Sales": table["Sales_2025"].map(fmt_metric_number),
        "Unique Visitors": table["UniqueVisitors_2025"].map(fmt_metric_number),
        visitor_col: table["Visitors YoY %"].map(lambda x: badge_span(x)),
        sales_col: table["Sales YoY %"].map(lambda x: badge_span(x)),
    })
    html_table(display, highlight_brands=[], sticky_brand=True)

def render_market_takeaways(data, market, selected_oems):
    yoy = yoy_table(data, market, selected_oems)
    if yoy.empty:
        return

    top_w2c = yoy.sort_values("ConversionPct_2025", ascending=False).iloc[0]
    top_visitors = yoy.sort_values("UniqueVisitors_2025", ascending=False).iloc[0]
    fastest_visitors = yoy.sort_values("Visitors YoY %", ascending=False).iloc[0]
    strongest_sales = yoy.sort_values("Sales YoY %", ascending=False).iloc[0]
    weakest_w2c_move = yoy.sort_values("Conv Var pp", ascending=True).iloc[0]

    cards = [
        ("Leader", f"{top_w2c['OEM']} leads W2C", f"{top_w2c['OEM']} records the highest Website-to-Contract Conversion Rate in {market}.", f"{top_w2c['ConversionPct_2025']:.2f}%", "positive"),
        ("Scale", f"{top_visitors['OEM']} leads visitor scale", f"{top_visitors['OEM']} has the largest unique visitor base in {market}.", fmt_metric_number(top_visitors["UniqueVisitors_2025"]), "primary"),
        ("Demand", f"{fastest_visitors['OEM']} has fastest visitor growth", f"Unique visitors grew {fmt_pct(fastest_visitors['Visitors YoY %'])} versus {PREVIOUS_LABEL}.", fmt_pct(fastest_visitors["Visitors YoY %"]), "positive"),
        ("Sales", f"{strongest_sales['OEM']} has strongest sales growth", f"Passenger car sales grew {fmt_pct(strongest_sales['Sales YoY %'])} versus {PREVIOUS_LABEL}.", fmt_pct(strongest_sales["Sales YoY %"]), "positive"),
        ("Risk", f"{weakest_w2c_move['OEM']} has weakest W2C movement", f"W2C moved {fmt_pp(weakest_w2c_move['Conv Var pp'])} versus {PREVIOUS_LABEL}.", fmt_pp(weakest_w2c_move["Conv Var pp"]), "risk"),
    ]

    html = "<div class='takeaway-grid'>"
    for label, title, copy, metric, style in cards:
        html += (
            f"<div class='takeaway-card {style}'>"
            f"<div class='takeaway-label'>{label}</div>"
            f"<div class='takeaway-title'>{title}</div>"
            f"<div class='takeaway-copy'>{copy}</div>"
            f"<div class='assistant-card-metric'>{metric}</div>"
            f"</div>"
        )
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)


def render_market_summary_page(data, market, selected_oems):
    section(f"Market Summary — {market}")
    st.caption(f"Executive view for {CURRENT_LABEL} vs {PREVIOUS_LABEL}: market growth, leading competitors and key takeaways.")
    if not render_market_kpi_cards(data, market, selected_oems):
        return

    section("Leading competitors — Top 10")
    top10_leaders_table(data, market, selected_oems)

    section("Five executive takeaways")
    render_market_takeaways(data, market, selected_oems)
    render_footer()

def render_bubble_side_table(df, market):
    latest = df[df["Year"] == 2025].copy()
    if latest.empty:
        latest = df.copy()

    latest = latest.sort_values("UniqueVisitors", ascending=False).head(15)

    rows = ""
    for _, r in latest.iterrows():
        cluster = cluster_for_oem(r["OEM"])
        color = CLUSTER_COLORS.get(cluster, "#6F6F6F")
        rows += (
            "<div class='bubble-row'>"
            f"<span class='bubble-row-dot' style='background:{color};'></span>"
            "<div>"
            f"<div class='bubble-row-brand'>{r['OEM']}</div>"
            f"<div class='bubble-row-meta'>{r['ConversionPct']:.2f}% W2C · {fmt_int(r['Sales'])} sales</div>"
            "</div>"
            f"<div class='bubble-row-value'>{fmt_short(r['UniqueVisitors'])}</div>"
            "</div>"
        )

    st.markdown(
        f"<div class='bubble-side-card'>"
        f"<div class='bubble-side-title'>{market} selected OEMs</div>"
        f"{rows}"
        f"</div>",
        unsafe_allow_html=True,
    )


def render_bubble_page(data, selected_oems, year_view=None, show_logos=False):
    section("OEM Bubble Chart")
    st.caption(f"Comparison selected: {CURRENT_LABEL} vs {PREVIOUS_LABEL}.")

    chart_col, control_col = st.columns([3, 1], gap="large")

    with control_col:
        st.markdown("### Bubble controls")
        local_year_view = st.selectbox("Bubble year view", ["Previous and current + shift", "Previous period", "Current period"], index=0)
        x_axis_metric = st.selectbox("X axis", ["Unique Visitors", "Passenger Car Sales"], index=0)

        default_market = summary_market if summary_market != "MM5" else "UK"
        selected_markets = st.multiselect(
            "Markets",
            ["UK", "France", "Germany", "Italy", "Spain"],
            default=[default_market],
            help="Default is one market to keep the view readable. Add more markets if you want cross-market comparison.",
        )

        category_options = available_oem_categories()
        default_categories = ["Volume Leaders"] if "Volume Leaders" in category_options else category_options[:1]
        selected_categories = st.multiselect(
            "OEM categories",
            category_options,
            default=default_categories,
            help="Select one or more OEM categories to include in the bubble chart.",
        )

        if selected_categories:
            default_oems = oems_for_categories(data, selected_categories)
        else:
            default_oems = sorted(data["OEM"].unique())

        oem_options = oems_for_categories(data, selected_categories)
        if not oem_options:
            oem_options = sorted(data["OEM"].unique())
        default_oems = [o for o in default_oems if o in oem_options] or oem_options
        local_oems = st.multiselect(
            "OEMs",
            oem_options,
            default=default_oems,
            help="Select the OEMs to plot.",
        )

    if not selected_markets:
        selected_markets = [summary_market if summary_market != "MM5" else "UK"]
    if not local_oems:
        local_oems = default_volume_oems(data)

    size_label = "Passenger Car Sales" if x_axis_metric == "Unique Visitors" else "Unique visitors"
    with chart_col:
        st.markdown(f"<div class='bubble-key'>Bubble size represents <b>{size_label}</b>.</div>", unsafe_allow_html=True)
        df = data[(data["Market"].isin(selected_markets)) & (data["OEM"].isin(local_oems))].copy()
        df["DisplayName"] = df["OEM"] + " · " + df["Market"]
        if local_year_view != "Previous and current + shift":
            selected_year = 2024 if local_year_view == "Previous period" else 2025
            df = df[df["Year"] == selected_year]
        if df.empty:
            st.info("No bubble chart data for this selection.")
        else:
            if len(selected_markets) > 1:
                df["Cluster"] = df["Market"]
            fig = build_bubble_chart(df, local_oems, "Selected markets", local_year_view, False, x_axis_metric=x_axis_metric, view_label="OEM")
            st.plotly_chart(fig, use_container_width=True)


def render_mm5_bubble_chart_page(data):
    section("MM5 Bubble Chart")
    st.caption(f"Comparison selected: {CURRENT_LABEL} vs {PREVIOUS_LABEL}. Compare market performance across MM5.")

    chart_col, control_col = st.columns([3, 1], gap="large")
    with control_col:
        st.markdown("### Bubble controls")
        local_year_view = st.selectbox("Bubble year view", ["Previous and current + shift", "Previous period", "Current period"], index=0, key="mm5_bubble_year")
        x_axis_metric = st.selectbox("X axis", ["Unique Visitors", "Passenger Car Sales"], index=0, key="mm5_bubble_x")

        category_options = available_oem_categories()
        selected_categories = st.multiselect(
            "OEM categories",
            category_options,
            default=category_options,
            key="mm5_bubble_categories",
            help="Select one or more OEM categories to aggregate into the market bubbles.",
        )

        if selected_categories:
            default_oems = oems_for_categories(data, selected_categories)
        else:
            default_oems = sorted(data["OEM"].unique())

        oem_options = oems_for_categories(data, selected_categories)
        if not oem_options:
            oem_options = sorted(data["OEM"].unique())
        default_oems = [o for o in default_oems if o in oem_options] or oem_options
        selected_oems_mm5 = st.multiselect(
            "OEMs",
            oem_options,
            default=default_oems,
            key="mm5_bubble_oems",
            help="The chart aggregates selected OEMs into the five market bubbles.",
        )

    if not selected_oems_mm5:
        selected_oems_mm5 = sorted(data["OEM"].unique())

    size_label = "Passenger Car Sales" if x_axis_metric == "Unique Visitors" else "Unique visitors"
    rows = []
    for market in ["UK", "France", "Germany", "Italy", "Spain"]:
        for year in [2024, 2025]:
            d = market_year(data, market, year, selected_oems_mm5)
            if d.empty:
                continue
            sales, visitors = d["Sales"].sum(), d["UniqueVisitors"].sum()
            rows.append({
                "OEM": market,
                "DisplayName": market,
                "Market": market,
                "Year": year,
                "Sales": sales,
                "UniqueVisitors": visitors,
                "ConversionPct": sales / visitors * 100 if visitors else 0,
                "Cluster": "MM5 Markets",
            })
    df = pd.DataFrame(rows)
    if local_year_view != "Previous and current + shift":
        selected_year = 2024 if local_year_view == "Previous period" else 2025
        df = df[df["Year"] == selected_year]

    with chart_col:
        st.markdown(f"<div class='bubble-key'>Bubble size represents <b>{size_label}</b>. Market bubbles are aggregated using the selected OEM/category/preset filter.</div>", unsafe_allow_html=True)
        if df.empty:
            st.info("No MM5 bubble chart data for this selection.")
        else:
            fig = build_bubble_chart(df, ["UK","France","Germany","Italy","Spain"], "MM5 markets", local_year_view, False, x_axis_metric=x_axis_metric, view_label="Market")
            st.plotly_chart(fig, use_container_width=True)
    render_footer()


def scorecard_table(data, market, selected_oems):
    yoy = yoy_table(data, market, selected_oems)
    if yoy.empty:
        return pd.DataFrame()
    out = yoy.copy()
    out["Category"] = out["OEM"].map(cluster_for_oem)
    out["W2C Ranking"] = out["ConversionPct_2025"].rank(method="first", ascending=False).astype(int)
    out["Sales Ranking"] = out["Sales_2025"].rank(method="first", ascending=False).astype(int)
    return out.sort_values("W2C Ranking")


def render_scorecard_page(data, selected_oems):
    section("Leadership scorecard")
    st.caption(f"Comparison selected: {CURRENT_LABEL} vs {PREVIOUS_LABEL}. Select any OEMs to compare; no brand is locked or pinned.")

    c1, c2, c3 = st.columns([1.15, 1.25, 0.85], gap="medium")
    with c1:
        market = st.selectbox("Scorecard market", MARKETS, index=0)
    with c2:
        sort_by = st.selectbox(
            "Sort table by",
            ["W2C ranking", "W2C rate", "W2C variance", "Passenger Car Sales", "Sales YoY", "Unique visitors", "Visitor YoY", "Visits to sale"],
            index=0,
        )
    with c3:
        descending = st.toggle("Sort large to small", value=True)

    score = scorecard_table(data, market, selected_oems)
    if score.empty:
        st.info("No scorecard data available.")
        return

    render_scorecard_podiums(score)

    sorted_score = sort_scorecard(score, sort_by, descending)
    visitor_col = f"Visitor YoY vs {PREVIOUS_LABEL}"
    sales_col = f"Sales YoY vs {PREVIOUS_LABEL}"
    w2c_var_col = f"W2C var vs {PREVIOUS_LABEL}"
    display = pd.DataFrame({
        "Brand": sorted_score["OEM"],
        "Category": sorted_score["Category"],
        "W2C rate": sorted_score["ConversionPct_2025"].map(lambda x: f"{x:.2f}%"),
        w2c_var_col: sorted_score["Conv Var pp"].map(lambda x: badge_span(x, suffix="pp")),
        "W2C ranking": sorted_score["W2C Ranking"].map(medal_for_rank),
        "Passenger Car Sales": sorted_score["Sales_2025"].map(fmt_metric_number),
        sales_col: sorted_score["Sales YoY %"].map(lambda x: badge_span(x)),
        "Unique visitors": sorted_score["UniqueVisitors_2025"].map(fmt_metric_number),
        visitor_col: sorted_score["Visitors YoY %"].map(lambda x: badge_span(x)),
        "Visits to sale": sorted_score["Visits to Sale 2025"].map(fmt_int),
    })
    html_table(display, highlight_brands=[], sticky_brand=True)
    st.download_button("Download scorecard as CSV", data=sorted_score.to_csv(index=False).encode("utf-8"), file_name=f"{market.lower()}_scorecard.csv", mime="text/csv")


# =========================
# App runtime
# =========================

render_hero()

raw_data = load_data()

st.sidebar.header("Comparison view")
comparison_view = st.sidebar.radio(
    "Choose the comparison period",
    list(COMPARISON_OPTIONS.keys()),
    index=0,
    help="Switch the full dashboard between annual 2024 vs 2025 and Jan-Apr 2025 vs Jan-Apr 2026.",
)

ACTIVE_COMPARISON = COMPARISON_OPTIONS[comparison_view]
PREVIOUS_LABEL = ACTIVE_COMPARISON["previous_label"]
CURRENT_LABEL = ACTIVE_COMPARISON["current_label"]

data = apply_comparison_view(raw_data, comparison_view)
all_oems = sorted(data["OEM"].unique())

# Startup sanity check. If this fails, the app should stop immediately with one useful message.
for m in MARKETS:
    y = yoy_table(data, m, None)
    if not y.empty:
        required_cols = {"Sales_2025", "Sales_2024", "UniqueVisitors_2025", "UniqueVisitors_2024", "Sales YoY %", "Visitors YoY %", "Conv Var pp"}
        missing = required_cols - set(y.columns)
        if missing:
            st.error(f"Data validation failed for {m}. Missing: {sorted(missing)}")
            st.stop()

st.sidebar.header("Filters")

page = st.sidebar.radio(
    "Dashboard page",
    ["Start Here", "Insight Report", "Use Cases", "Market Summary", "OEM Bubble Chart", "MM5 Bubble Chart", "Scorecard", "Data Assistant"],
    index=0,
)

summary_market = st.sidebar.selectbox("Summary market", MARKETS, index=0)


cluster_options = ["All categories"] + list(OEM_CLUSTERS.keys())

selected_clusters = st.sidebar.multiselect(
    "OEM categories",
    cluster_options,
    default=["All categories"],
    help="Filter the dashboard by OEM grouping. These categories are based on the brand clusters used in the companion marketing dashboards.",
)

cluster_filtered_oems = oems_for_clusters(selected_clusters, all_oems)

default_oems = cluster_filtered_oems

selected_oems = st.sidebar.multiselect(
    "OEMs",
    cluster_filtered_oems,
    default=default_oems,
)

show_logos = False

selected_oems = selected_or_all(selected_oems, cluster_filtered_oems)

if page == "Start Here":
    render_start_here_page(data)
elif page == "Insight Report":
    render_insight_report_page(data)
elif page == "Use Cases":
    render_use_cases_page(data)
elif page == "Market Summary":
    render_market_summary_page(data, summary_market, selected_oems)
elif page == "OEM Bubble Chart":
    render_bubble_page(data, selected_oems, show_logos=False)
elif page == "MM5 Bubble Chart":
    render_mm5_bubble_chart_page(data)
elif page == "Scorecard":
    render_scorecard_page(data, selected_oems)
else:
    render_data_assistant_page(data, selected_oems)
